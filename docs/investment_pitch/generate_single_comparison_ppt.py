#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成单页传统机器人 vs Brain 对比PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_single_comparison_slide():
    """创建单页对比PPT"""
    prs = Presentation()

    # 设置16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ========== 标题区域 ==========
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "传统机器人 vs Brain 系统"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # ========== 左侧：传统机器人 ==========
    # 主框架
    trad_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.0), Inches(4.5), Inches(6.0)
    )
    trad_box.fill.solid()
    trad_box.fill.fore_color.rgb = RGBColor(255, 235, 238)  # 浅红色
    trad_box.line.color.rgb = RGBColor(229, 57, 53)
    trad_box.line.width = Pt(4)

    # 标题
    trad_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.15), Inches(4.1), Inches(0.5)
    )
    tf = trad_title.text_frame
    tf.text = "❌ 传统机器人"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(183, 28, 28)
    p.alignment = PP_ALIGN.CENTER

    # 5个核心痛点
    pain_points = [
        ("🔴 开发周期", "6-12个月/平台", "重复造轮子"),
        ("🔴 智能程度", "预设程序", "环境变化就傻"),
        ("🔴 协作能力", "单机作战", "无法多机协同"),
        ("🔴 适应性", "硬编码", "需重新编程"),
        ("💰 成本", "50-100万/年", "维护困难")
    ]

    y_start = Inches(1.85)
    for icon, title, desc in pain_points:
        # 痛点框
        pain_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), y_start, Inches(4.1), Inches(0.85)
        )
        pain_box.fill.solid()
        pain_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        pain_box.line.color.rgb = RGBColor(229, 57, 53)
        pain_box.line.width = Pt(2)

        # 标题
        title_text = slide.shapes.add_textbox(
            Inches(0.65), y_start + Inches(0.08), Inches(3.8), Inches(0.3)
        )
        tf = title_text.text_frame
        tf.text = f"{title}: {desc}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(183, 28, 28)

        # 详细说明
        desc_text = slide.shapes.add_textbox(
            Inches(0.65), y_start + Inches(0.4), Inches(3.8), Inches(0.35)
        )
        tf = desc_text.text_frame
        tf.word_wrap = True
        tf.text = icon.replace('🔴', '•').replace('💰', '•')
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(100, 100, 100)

        y_start += Inches(0.95)

    # ========== 右侧：Brain系统 ==========
    # 主框架
    brain_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.0), Inches(4.5), Inches(6.0)
    )
    brain_box.fill.solid()
    brain_box.fill.fore_color.rgb = RGBColor(232, 245, 233)  # 浅绿色
    brain_box.line.color.rgb = RGBColor(67, 160, 71)
    brain_box.line.width = Pt(4)

    # 标题
    brain_title = slide.shapes.add_textbox(
        Inches(5.4), Inches(1.15), Inches(4.1), Inches(0.5)
    )
    tf = brain_title.text_frame
    tf.text = "✅ Brain 系统"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 94, 32)
    p.alignment = PP_ALIGN.CENTER

    # 5个核心优势（与左侧对应）
    advantages = [
        ("🟢 开发周期", "3个月全平台", "代码复用90%"),
        ("🟢 智能程度", "AI理解决策", "自适应环境变化"),
        ("🟢 协作能力", "多机协同", "效率提升3-5倍"),
        ("🟢 适应性", "自然语言", "自动重规划"),
        ("💰 成本", "15-30万/年", "易维护扩展")
    ]

    y_start = Inches(1.85)
    for icon, title, desc in advantages:
        # 优势框
        adv_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.4), y_start, Inches(4.1), Inches(0.85)
        )
        adv_box.fill.solid()
        adv_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        adv_box.line.color.rgb = RGBColor(67, 160, 71)
        adv_box.line.width = Pt(2)

        # 标题
        title_text = slide.shapes.add_textbox(
            Inches(5.55), y_start + Inches(0.08), Inches(3.8), Inches(0.3)
        )
        tf = title_text.text_frame
        tf.text = f"{title}: {desc}"
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(27, 94, 32)

        # 详细说明
        desc_text = slide.shapes.add_textbox(
            Inches(5.55), y_start + Inches(0.4), Inches(3.8), Inches(0.35)
        )
        tf = desc_text.text_frame
        tf.word_wrap = True
        tf.text = icon.replace('🟢', '•').replace('💰', '•')
        p = tf.paragraphs[0]
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(100, 100, 100)

        y_start += Inches(0.95)

    # ========== 中间VS箭头 ==========
    vs_box = slide.shapes.add_textbox(
        Inches(4.5), Inches(3.5), Inches(1), Inches(0.6)
    )
    tf = vs_box.text_frame
    tf.text = "VS"
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 152, 0)
    p.alignment = PP_ALIGN.CENTER

    # ========== 底部核心优势 ==========
    # 核心技术框
    tech_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(7.15), Inches(9.4), Inches(0.25)
    )
    tech_box.fill.solid()
    tech_box.fill.fore_color.rgb = RGBColor(33, 150, 243)  # 蓝色
    tech_box.line.color.rgb = RGBColor(13, 71, 161)
    tech_box.line.width = Pt(2)

    # 核心技术说明
    tech_text = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.18), Inches(9), Inches(0.2)
    )
    tf = tech_text.text_frame
    tf.word_wrap = True
    tf.text = "💡 核心技术：World Model 理解环境  +  CoT 推理决策  +  HTN 分层规划  +  自适应执行  =  通用智能机器人操作系统"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 保存
    output_path = "/media/yangyuhui/CODES1/Brain/docs/investment_pitch/Brain_Comparison_Single.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页幻灯片")
    print()
    print("🎯 设计特点：")
    print("  • 左右对比：红色（传统痛点）vs 绿色（Brain优势）")
    print("  • 5个核心维度：开发周期、智能程度、协作能力、适应性、成本")
    print("  • 数据支撑：每个对比都有具体数字")
    print("  • 视觉冲击：用颜色和图标强调对比")
    print("  • 底部总结：突出4大核心技术")


if __name__ == "__main__":
    create_single_comparison_slide()
