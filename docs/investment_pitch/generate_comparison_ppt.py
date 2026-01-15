#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成传统机器人 vs Brain 对比PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_comparison_slide(prs):
    """创建传统 vs Brain 对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "传统机器人 vs Brain 系统"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.75), Inches(9.4), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "为什么需要 Brain？传统机器人的核心痛点 vs 我们的解决方案"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # 左右对比框架
    # ========== 传统机器人（左侧）==========
    traditional_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.3), Inches(4.5), Inches(5.8)
    )
    traditional_box.fill.solid()
    traditional_box.fill.fore_color.rgb = RGBColor(255, 235, 238)  # 浅红色背景
    traditional_box.line.color.rgb = RGBColor(229, 57, 53)  # 红色边框
    traditional_box.line.width = Pt(3)

    # 传统机器人标题
    trad_title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.45), Inches(4.1), Inches(0.5)
    )
    trad_title_frame = trad_title_box.text_frame
    trad_title_frame.text = "❌ 传统机器人系统"
    p = trad_title_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(183, 28, 28)
    p.alignment = PP_ALIGN.CENTER

    # 传统机器人内容
    trad_items = [
        ("开发成本高", [
            "• 每个平台独立开发: 6-12个月",
            "• 代码无法复用: 重复造轮子",
            "• 人力成本: 50-100万元/年",
            "• 维护困难: 多套代码并行维护"
        ]),
        ("智能程度低", [
            "• 只能执行预设任务: 编程写死",
            "• 无法理解环境: 盲目执行",
            "• 遇到变化就傻: 停机待命",
            "• 依赖人工干预: 远程遥控"
        ]),
        ("协作能力差", [
            "• 单机作战: 无法协同",
            "• 通信困难: 协议不统一",
            "• 任务分配: 人工指定",
            "• 效率低下: 重复劳动"
        ]),
        ("适应性弱", [
            "• 环境变化: 需重新编程",
            "• 任务变化: 需重新部署",
            "• 失败处理: 等待人工",
            "• 扩展困难: 硬编码限制"
        ])
    ]

    y_start = Inches(2.0)
    for i, (title, items) in enumerate(trad_items):
        # 标题
        item_title = slide.shapes.add_textbox(
            Inches(0.5), y_start, Inches(4.1), Inches(0.3)
        )
        tf = item_title.text_frame
        tf.text = f"🔴 {title}"
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(183, 28, 28)

        # 内容列表
        item_content = slide.shapes.add_textbox(
            Inches(0.5), y_start + Inches(0.28), Inches(4.1), Inches(1.0)
        )
        tf = item_content.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.space_before = Pt(2)
            p.space_after = Pt(1)

        y_start += Inches(1.2)

    # ========== Brain系统（右侧）==========
    brain_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.3), Inches(4.5), Inches(5.8)
    )
    brain_box.fill.solid()
    brain_box.fill.fore_color.rgb = RGBColor(232, 245, 233)  # 浅绿色背景
    brain_box.line.color.rgb = RGBColor(67, 160, 71)  # 绿色边框
    brain_box.line.width = Pt(3)

    # Brain标题
    brain_title_box = slide.shapes.add_textbox(
        Inches(5.4), Inches(1.45), Inches(4.1), Inches(0.5)
    )
    brain_title_frame = brain_title_box.text_frame
    brain_title_frame.text = "✅ Brain 智能系统"
    p = brain_title_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 94, 32)
    p.alignment = PP_ALIGN.CENTER

    # Brain内容（与传统对应）
    brain_items = [
        ("开发成本低", [
            "• 一次开发多平台复用: 3个月内",
            "• 代码复用率90%: 大幅降低成本",
            "• 人力成本: 15-30万元/年",
            "• 易于维护: 统一架构"
        ]),
        ("智能程度高", [
            "• VLM理解场景: 看懂环境",
            "• LLM推理决策: 自主规划",
            "• 遇到变化自适应: 实时调整",
            "• 自主执行: 无需人工干预"
        ]),
        ("协作能力强", [
            "• 多机协同: 自动任务分配",
            "• 统一通信: 标准化接口",
            "• 智能调度: 动态优化",
            "• 效率提升: 3-5倍"
        ]),
        ("适应性极强", [
            "• 环境变化: 自动重规划",
            "• 任务变化: 自然语言理解",
            "• 失败处理: 自动恢复",
            "• 灵活扩展: 插件化架构"
        ])
    ]

    y_start = Inches(2.0)
    for i, (title, items) in enumerate(brain_items):
        # 标题
        item_title = slide.shapes.add_textbox(
            Inches(5.4), y_start, Inches(4.1), Inches(0.3)
        )
        tf = item_title.text_frame
        tf.text = f"🟢 {title}"
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(27, 94, 32)

        # 内容列表
        item_content = slide.shapes.add_textbox(
            Inches(5.4), y_start + Inches(0.28), Inches(4.1), Inches(1.0)
        )
        tf = item_content.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.space_before = Pt(2)
            p.space_after = Pt(1)

        y_start += Inches(1.2)

    # 中间VS箭头
    arrow_box = slide.shapes.add_textbox(
        Inches(4.6), Inches(3.8), Inches(0.8), Inches(0.6)
    )
    arrow_frame = arrow_box.text_frame
    arrow_frame.text = "VS"
    p = arrow_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 152, 0)
    p.alignment = PP_ALIGN.CENTER

    # 底部总结框
    summary_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(7.2), Inches(9.4), Inches(0.6)
    )
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = RGBColor(33, 150, 243)  # 蓝色
    summary_box.line.color.rgb = RGBColor(13, 71, 161)
    summary_box.line.width = Pt(2)

    summary_text = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.25), Inches(9), Inches(0.5)
    )
    tf = summary_text.text_frame
    tf.word_wrap = True
    tf.text = "💡 核心优势：World Model 理解世界 + CoT 推理决策 + HTN 分层规划 + 自适应执行 = 通用智能机器人操作系统"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def create_metrics_slide(prs):
    """创建量化对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "量化对比：数据说话"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 6个指标卡片
    metrics = [
        ("开发周期", "6-12个月", "3个月", "70%", RGBColor(255, 152, 0)),
        ("代码复用率", "0-10%", "90%+", "9倍", RGBColor(156, 39, 176)),
        ("开发成本", "50-100万", "15-30万", "70%", RGBColor(46, 125, 50)),
        ("任务适应性", "需重新编程", "自然语言理解", "质变", RGBColor(0, 151, 167)),
        ("多机协作", "不支持", "原生支持", "从0到1", RGBColor(233, 30, 99)),
        ("失败恢复", "人工处理", "自动恢复", "100%", RGBColor(255, 193, 7))
    ]

    # 2行3列布局
    card_width = Inches(3)
    card_height = Inches(2.2)
    card_spacing_x = Inches(0.15)
    card_spacing_y = Inches(0.15)
    start_x = Inches(0.4)
    start_y = Inches(1.0)

    for i, (metric, trad, brain, improvement, color) in enumerate(metrics):
        row = i // 3
        col = i % 3

        x = start_x + col * (card_width + card_spacing_x)
        y = start_y + row * (card_height + card_spacing_y)

        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(250, 250, 250)
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # 指标名称
        metric_box = slide.shapes.add_textbox(
            x + Inches(0.1), y + Inches(0.1), card_width - Inches(0.2), Inches(0.35)
        )
        tf = metric_box.text_frame
        tf.text = metric
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.CENTER

        # 传统系统
        trad_box = slide.shapes.add_textbox(
            x + Inches(0.1), y + Inches(0.5), card_width - Inches(0.2), Inches(0.3)
        )
        tf = trad_box.text_frame
        tf.text = f"传统: {trad}"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(183, 28, 28)

        # Brain系统
        brain_box = slide.shapes.add_textbox(
            x + Inches(0.1), y + Inches(0.8), card_width - Inches(0.2), Inches(0.3)
        )
        tf = brain_box.text_frame
        tf.text = f"Brain: {brain}"
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(27, 94, 32)

        # 改进幅度
        improve_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x + Inches(0.3), y + Inches(1.15), card_width - Inches(0.6), Inches(0.4)
        )
        improve_box.fill.solid()
        improve_box.fill.fore_color.rgb = color

        improve_text = slide.shapes.add_textbox(
            x + Inches(0.3), y + Inches(1.2), card_width - Inches(0.6), Inches(0.3)
        )
        tf = improve_text.text_frame
        tf.text = f"↑ {improvement}"
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # 底部总结
    summary_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.8)
    )
    tf = summary_box.text_frame
    tf.word_wrap = True
    tf.text = """🎯 关键结论：Brain 系统在所有核心指标上都有显著提升，特别是代码复用率（9倍）和开发成本（降低70%）
这意味着：更快的上市时间、更低的开发成本、更强的智能能力、更好的扩展性"""
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.5


def create_architecture_comparison_slide(prs):
    """创建架构对比页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "架构对比：硬编码 vs 智能化"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 左侧：传统架构
    trad_arch_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(1.0), Inches(4.5), Inches(5.5)
    )
    trad_arch_box.fill.solid()
    trad_arch_box.fill.fore_color.rgb = RGBColor(255, 235, 238)
    trad_arch_box.line.color.rgb = RGBColor(229, 57, 53)
    trad_arch_box.line.width = Pt(2)

    # 传统架构标题
    trad_arch_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.15), Inches(4.1), Inches(0.4)
    )
    tf = trad_arch_title.text_frame
    tf.text = "传统架构：硬编码规则"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(183, 28, 28)
    p.alignment = PP_ALIGN.CENTER

    # 传统架构层
    trad_layers = [
        ("应用层", "if position == kitchen:\n    goto(kitchen)\nif object == cup:\n    grasp(cup)\n# 硬编码逻辑"),
        ("控制层", "motor_control(speed, angle)\n# 底层控制"),
        ("驱动层", "PWM输出\n传感器读取")
    ]

    y_start = Inches(1.7)
    for name, desc in trad_layers:
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), y_start, Inches(4.1), Inches(1.1)
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        layer_box.line.color.rgb = RGBColor(229, 57, 53)
        layer_box.line.width = Pt(1.5)

        # 层名称
        name_box = slide.shapes.add_textbox(
            Inches(0.65), y_start + Inches(0.1), Inches(1), Inches(0.3)
        )
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(183, 28, 28)

        # 描述
        desc_box = slide.shapes.add_textbox(
            Inches(0.5), y_start + Inches(0.4), Inches(4.1), Inches(0.6)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.family = "Courier New"
            paragraph.font.color.rgb = RGBColor(80, 80, 80)

        y_start += Inches(1.3)

    # 传统架构问题说明
    trad_problem = slide.shapes.add_textbox(
        Inches(0.5), y_start + Inches(0.2), Inches(4.1), Inches(0.8)
    )
    tf = trad_problem.text_frame
    tf.word_wrap = True
    tf.text = "❌ 问题：\n• 环境变化需要重写代码\n• 任务增加需要添加if-else\n• 无法理解复杂场景\n• 维护成本高"
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(183, 28, 28)

    # 右侧：Brain架构
    brain_arch_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.0), Inches(4.5), Inches(5.5)
    )
    brain_arch_box.fill.solid()
    brain_arch_box.fill.fore_color.rgb = RGBColor(232, 245, 233)
    brain_arch_box.line.color.rgb = RGBColor(67, 160, 71)
    brain_arch_box.line.width = Pt(2)

    # Brain架构标题
    brain_arch_title = slide.shapes.add_textbox(
        Inches(5.4), Inches(1.15), Inches(4.1), Inches(0.4)
    )
    tf = brain_arch_title.text_frame
    tf.text = "Brain 架构：AI 驱动"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(27, 94, 32)
    p.alignment = PP_ALIGN.CENTER

    # Brain架构层
    brain_layers = [
        ("执行层", "自适应执行引擎\n失败自动恢复\n实时监控调整"),
        ("规划层", "HTN分层规划\n动态推理\n任务自动分解"),
        ("认知层", "World Model + CoT\n理解环境\n推理决策"),
        ("感知层", "VLM多模态融合\n看懂场景\n语义理解")
    ]

    y_start = Inches(1.7)
    for name, desc in brain_layers:
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.4), y_start, Inches(4.1), Inches(1.1)
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        layer_box.line.color.rgb = RGBColor(67, 160, 71)
        layer_box.line.width = Pt(1.5)

        # 层名称
        name_box = slide.shapes.add_textbox(
            Inches(5.55), y_start + Inches(0.1), Inches(1), Inches(0.3)
        )
        tf = name_box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(27, 94, 32)

        # 描述
        desc_box = slide.shapes.add_textbox(
            Inches(6.7), y_start + Inches(0.15), Inches(2.7), Inches(0.8)
        )
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.text = desc
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(80, 80, 80)
            paragraph.line_spacing = 1.3

        y_start += Inches(1.3)

    # Brain架构优势说明
    brain_advantage = slide.shapes.add_textbox(
        Inches(5.4), y_start + Inches(0.2), Inches(4.1), Inches(0.8)
    )
    tf = brain_advantage.text_frame
    tf.word_wrap = True
    tf.text = "✅ 优势：\n• 自然语言理解任务\n• 环境变化自动适应\n• 智能推理决策\n• 一次开发多平台复用"
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(27, 94, 32)

    # 中间箭头
    arrow_text = slide.shapes.add_textbox(
        Inches(4.5), Inches(3.5), Inches(1), Inches(0.5)
    )
    tf = arrow_text.text_frame
    tf.text = "进化 →"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 152, 0)
    p.alignment = PP_ALIGN.CENTER

    # 底部说明
    summary_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.7), Inches(9), Inches(0.6)
    )
    tf = summary_box.text_frame
    tf.word_wrap = True
    tf.text = "🚀 核心差异：传统系统是 '程序控制'（写死的逻辑），Brain 是 '智能决策'（AI驱动的自适应）"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER


def main():
    """生成对比PPT"""
    prs = Presentation()

    # 设置16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 创建3页
    create_comparison_slide(prs)
    create_metrics_slide(prs)
    create_architecture_comparison_slide(prs)

    # 保存
    output_path = "/media/yangyuhui/CODES1/Brain/docs/investment_pitch/Brain_Comparison.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页幻灯片")
    print()
    print("📋 内容：")
    print("  第1页：传统 vs Brain - 核心痛点对比（开发成本/智能程度/协作能力/适应性）")
    print("  第2页：量化对比 - 6个关键指标数据对比（开发周期/代码复用率/成本等）")
    print("  第3页：架构对比 - 硬编码规则 vs AI驱动架构")


if __name__ == "__main__":
    main()
