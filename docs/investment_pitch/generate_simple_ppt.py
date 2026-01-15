#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成Brain项目投资路演PPT - 精简版（3页）
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def add_architecture_diagram(slide, left, top, width, height):
    """添加架构图"""
    # 定义颜色
    colors = {
        'perception': RGBColor(255, 200, 87),    # 黄色
        'cognitive': RGBColor(144, 202, 249),    # 蓝色
        'planning': RGBColor(129, 199, 132),     # 绿色
        'execution': RGBColor(165, 214, 167),    # 浅绿
        'platform': RGBColor(230, 124, 115)      # 红色
    }

    # 感知层
    perception_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, Inches(0.6)
    )
    perception_box.fill.solid()
    perception_box.fill.fore_color.rgb = colors['perception']
    perception_box.line.color.rgb = RGBColor(0, 0, 0)
    perception_box.line.width = Pt(2)

    tf = perception_box.text_frame
    tf.word_wrap = True
    tf.text = "Perception Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 认知层
    cognitive_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top + Inches(0.8), width, Inches(0.6)
    )
    cognitive_box.fill.solid()
    cognitive_box.fill.fore_color.rgb = colors['cognitive']
    cognitive_box.line.color.rgb = RGBColor(0, 0, 0)
    cognitive_box.line.width = Pt(2)

    tf = cognitive_box.text_frame
    tf.word_wrap = True
    tf.text = "Cognitive Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 规划层
    planning_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top + Inches(1.6), width, Inches(0.6)
    )
    planning_box.fill.solid()
    planning_box.fill.fore_color.rgb = colors['planning']
    planning_box.line.color.rgb = RGBColor(0, 0, 0)
    planning_box.line.width = Pt(2)

    tf = planning_box.text_frame
    tf.word_wrap = True
    tf.text = "Planning Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 执行层
    execution_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top + Inches(2.4), width, Inches(0.6)
    )
    execution_box.fill.solid()
    execution_box.fill.fore_color.rgb = colors['execution']
    execution_box.line.color.rgb = RGBColor(0, 0, 0)
    execution_box.line.width = Pt(2)

    tf = execution_box.text_frame
    tf.word_wrap = True
    tf.text = "Execution Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 添加箭头连接
    for i in range(3):
        y = top + Inches(0.6) + i * Inches(0.8)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            left + Inches(1.5), y, Inches(0.3), Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
        arrow.line.fill.background()

    return top + Inches(3.2)


def create_slide_1_why(prs):
    """第1页：为什么做这个项目"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Why Brain? Market Opportunity"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)  # 深蓝
    p.alignment = PP_ALIGN.CENTER

    # 左侧：市场痛点
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(4.2), Inches(5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    p.text = "Market Pain Points"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)  # 红色
    p.space_after = Pt(20)

    pains = [
        "High Development Cost",
        "6-12 months per platform",
        "No code reuse",
        "",
        "Low Intelligence",
        "Pre-programmed tasks only",
        "No environment understanding",
        "",
        "Poor Collaboration",
        "Platforms cannot work together",
        "Manual coordination needed"
    ]

    for pain in pains:
        p = left_frame.add_paragraph()
        p.text = pain
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(8)

    # 右侧：我们的机会
    right_box = slide.shapes.add_textbox(
        Inches(5.3), Inches(1.5), Inches(4.2), Inches(5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    p = right_frame.paragraphs[0]
    p.text = "Our Opportunity"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)  # 蓝色
    p.space_after = Pt(20)

    opportunities = [
        "300B RMB Market by 2025",
        "Small/Mini/Micro Robots",
        "",
        "Government Support",
        "14th Five-Year Plan priority",
        "",
        "Technology Ready",
        "AI + Edge Computing + 5G",
        "",
        "Open Market",
        "No universal OS exists yet"
    ]

    for opp in opportunities:
        p = right_frame.add_paragraph()
        p.text = opp
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(8)


def create_slide_2_how(prs):
    """第2页：准备怎么做（技术架构）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "How? Technical Architecture"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 左侧：架构图
    arch_left = Inches(0.8)
    arch_top = Inches(1.5)
    arch_width = Inches(3.2)

    y_after_arch = add_architecture_diagram(
        slide, arch_left, arch_top, arch_width, Inches(3)
    )

    # 添加平台
    platform_y = y_after_arch + Inches(0.2)

    platforms = ["Drone", "UGV", "USV"]
    platform_width = Inches(0.9)
    platform_spacing = Inches(0.15)
    total_platform_width = len(platforms) * platform_width + (len(platforms) - 1) * platform_spacing
    platform_left = arch_left + (arch_width - total_platform_width) / 2

    for i, platform in enumerate(platforms):
        platform_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            platform_left + i * (platform_width + platform_spacing),
            platform_y,
            platform_width,
            Inches(0.5)
        )
        platform_box.fill.solid()
        platform_box.fill.fore_color.rgb = RGBColor(255, 152, 0)  # 橙色
        platform_box.line.color.rgb = RGBColor(0, 0, 0)
        platform_box.line.width = Pt(2)

        tf = platform_box.text_frame
        tf.word_wrap = True
        tf.text = platform
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # 右侧：核心能力说明
    right_box = slide.shapes.add_textbox(
        Inches(4.5), Inches(1.5), Inches(5), Inches(5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    p = right_frame.paragraphs[0]
    p.text = "Core Capabilities"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_after = Pt(16)

    capabilities = [
        ("Perception", "Multi-sensor fusion\nLiDAR + Camera + IMU\nSemantic understanding"),
        ("Cognitive", "World Model\nEnvironment reasoning\nChain-of-Thought"),
        ("Planning", "HTN task decomposition\nDynamic insertion\nFailure recovery"),
        ("Execution", "Adaptive executor\nReal-time monitoring\nAuto-replanning"),
        ("", ""),
        ("Key Metric", "168 test cases, 100% pass\nFull-stack self-developed\nOpen architecture")
    ]

    for title, desc in capabilities:
        if title:
            p = right_frame.add_paragraph()
            p.text = f"{title}:"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.space_before = Pt(8)
            p.space_after = Pt(2)

            p = right_frame.add_paragraph()
            p.text = desc
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_after = Pt(8)
        else:
            p = right_frame.add_paragraph()
            p.text = ""
            p.space_after = Pt(4)


def create_slide_3_applications(prs):
    """第3页：做出来用在哪"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Applications: Where & What"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 三列：无人机、无人车、无人船
    columns = [
        ("Drone Applications", [
            "Emergency Rescue",
            "Search survivors, Auto-drop supplies",
            "10x efficiency improvement",
            "",
            "Power Line Inspection",
            "AI defect detection",
            "Cost: 1/5 of manual"
        ]),
        ("UGV Applications", [
            "Warehouse Logistics",
            "Multi-vehicle coordination",
            "3x efficiency, 24/7 operation",
            "",
            "Security Patrol",
            "Face recognition, Anomaly detection",
            "70% cost reduction"
        ]),
        ("USV Applications", [
            "Water Quality Monitor",
            "Autonomous sampling, Real-time test",
            "80% cost reduction",
            "",
            "Underwater Inspection",
            "Sonar + Vision dual-mode",
            "Bridge structure detection"
        ])
    ]

    x_positions = [Inches(0.5), Inches(3.5), Inches(6.5)]
    y_start = Inches(1.5)

    for i, (title, items) in enumerate(columns):
        # 列标题
        col_box = slide.shapes.add_textbox(
            x_positions[i], y_start, Inches(3), Inches(0.6)
        )
        col_frame = col_box.text_frame
        col_frame.word_wrap = True
        col_frame.text = title
        p = col_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.alignment = PP_ALIGN.CENTER

        # 内容
        content_box = slide.shapes.add_textbox(
            x_positions[i], y_start + Inches(0.7), Inches(3), Inches(4.5)
        )
        content_frame = content_box.text_frame
        content_frame.word_wrap = True

        for item in items:
            p = content_frame.add_paragraph() if content_frame.paragraphs[0].text else content_frame.paragraphs[0]
            p.text = item
            if item == "":
                p.space_after = Pt(8)
            else:
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.space_before = Pt(4)
                p.space_after = Pt(2)

    # 底部：商业价值
    business_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.5), Inches(9), Inches(0.8)
    )
    business_frame = business_box.text_frame
    business_frame.word_wrap = True
    business_frame.text = "Business Model: B2B License (5K-50K RMB/robot/year) | Government Projects (1M-10M RMB) | Developer Ecosystem"
    p = business_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)
    p.alignment = PP_ALIGN.CENTER


def main():
    """生成精简版PPT"""
    prs = Presentation()

    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 创建3页幻灯片
    create_slide_1_why(prs)
    create_slide_2_how(prs)
    create_slide_3_applications(prs)

    # 保存PPT
    output_path = "/media/yangyuhui/CODES1/Brain/docs/investment_pitch/Brain_Investment_Pitch_Simple.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页幻灯片")


if __name__ == "__main__":
    main()
