#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成Brain项目完整架构图PPT - 专业美观版
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement


def add_connector(slide, start_shape, end_shape, x1, y1, x2, y2, width=Pt(1.5), color=RGBColor(80, 80, 80)):
    """添加连接线"""
    left = min(x1, x2)
    top = min(y1, y2)
    width_shape = abs(x2 - x1)
    height_shape = abs(y2 - y1)

    line = slide.shapes.add_shape(
        MSO_SHAPE.LINE,
        left, top, width_shape, height_shape
    )
    line.line.color.rgb = color
    line.line.width = width
    return line


def create_title_slide(prs):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Brain 系统完整技术架构"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.4), Inches(9), Inches(0.6)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "感知层 → 认知层 → 规划层 → 执行层"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # 三大技术优势
    advantages = [
        "World Model 驱动的智能理解",
        "HTN 分层规划 + 动态推理",
        "自适应执行引擎"
    ]

    y_start = Inches(2.8)
    for i, adv in enumerate(advantages):
        adv_box = slide.shapes.add_textbox(
            Inches(1.5), y_start + i * Inches(0.7), Inches(7), Inches(0.5)
        )
        adv_frame = adv_box.text_frame
        adv_frame.text = f"● {adv}"
        p = adv_frame.paragraphs[0]
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.space_before = Pt(8)

    # 底部说明
    note_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5), Inches(9), Inches(1.5)
    )
    note_frame = note_box.text_frame
    note_frame.word_wrap = True
    note_frame.text = """核心技术特点：
• VLM 视觉语言模型（LLaVA/MiniCPM-V）实现场景理解与目标搜索
• LLM 大语言模型（GPT-4/Claude）实现 Chain-of-Thought 推理
• World Model 世界模型提供几何、语义、动态三维态势
• HTN 分层规划支持任务层→技能层→动作层分解
• 动态推理实现运行时插入前置条件和自动重规划
• 自适应执行引擎实时监控并自动恢复失败"""
    p = note_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.line_spacing = 1.5


def create_architecture_slide(prs):
    """创建完整架构图页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Brain 系统完整架构图"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # ========== 传感器层 ==========
    sensor_y = Inches(0.8)
    sensor_width = Inches(1.7)
    sensor_height = Inches(0.5)
    sensor_spacing = Inches(0.15)
    sensor_start_x = Inches(0.5)

    sensors = [
        ("LiDAR\n激光雷达", "#FF9800"),
        ("摄像头\nRGB/深度/热成像", "#FF9800"),
        ("IMU\n惯性测量", "#FF9800"),
        ("GPS\n定位", "#FF9800")
    ]

    sensor_boxes = []
    for i, (text, color) in enumerate(sensors):
        x = sensor_start_x + i * (sensor_width + sensor_spacing)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, sensor_y, sensor_width, sensor_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 224, 178)
        box.line.color.rgb = RGBColor(255, 152, 0)
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.word_wrap = True
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        sensor_boxes.append((box, x, sensor_y, sensor_width, sensor_height))

    # ========== 感知层 ==========
    perception_y = sensor_y + Inches(0.7)
    perception_height = Inches(1.1)

    perception_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), perception_y, Inches(9.4), perception_height
    )
    perception_box.fill.solid()
    perception_box.fill.fore_color.rgb = RGBColor(255, 249, 196)
    perception_box.line.color.rgb = RGBColor(251, 192, 45)
    perception_box.line.width = Pt(2.5)

    # 感知层标题
    perception_title = slide.shapes.add_textbox(
        Inches(0.45), perception_y + Inches(0.05), Inches(2), Inches(0.3)
    )
    tf = perception_title.text_frame
    tf.text = "感知层 Perception Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 120, 0)

    # 感知层模块
    perception_modules = [
        ("SensorInput\n传感器输入", Inches(0.5)),
        ("PointCloud\n点云处理", Inches(2.15)),
        ("ObjectDetector\nYOLO检测", Inches(3.8)),
        ("VLM视觉模型\nLLaVA/MiniCPM", Inches(5.45)),
        ("FusionEngine\n融合引擎", Inches(7.1)),
        ("SituationalMap\n态势图", Inches(8.3))
    ]

    perception_module_width = Inches(1.4)
    perception_module_height = Inches(0.7)
    perception_module_y = perception_y + Inches(0.35)

    for name, x in perception_modules:
        # VLM用红色突出
        if "VLM" in name or "视觉模型" in name:
            bg_color = RGBColor(255, 235, 238)
            border_color = RGBColor(229, 57, 53)
            border_width = Pt(3)
        else:
            bg_color = RGBColor(255, 245, 157)
            border_color = RGBColor(251, 192, 45)
            border_width = Pt(1.5)

        module = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, perception_module_y, perception_module_width, perception_module_height
        )
        module.fill.solid()
        module.fill.fore_color.rgb = bg_color
        module.line.color.rgb = border_color
        module.line.width = border_width

        tf = module.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # ========== 认知层 ==========
    cognitive_y = perception_y + Inches(1.3)
    cognitive_height = Inches(1.1)

    cognitive_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), cognitive_y, Inches(9.4), cognitive_height
    )
    cognitive_box.fill.solid()
    cognitive_box.fill.fore_color.rgb = RGBColor(179, 229, 252)
    cognitive_box.line.color.rgb = RGBColor(2, 136, 209)
    cognitive_box.line.width = Pt(2.5)

    # 认知层标题
    cognitive_title = slide.shapes.add_textbox(
        Inches(0.45), cognitive_y + Inches(0.05), Inches(2.5), Inches(0.3)
    )
    tf = cognitive_title.text_frame
    tf.text = "认知层 Cognitive Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 80, 150)

    # 认知层模块
    cognitive_modules = [
        ("PerceptionParser\n感知解析", Inches(0.5)),
        ("World Model\n世界模型", Inches(2.0)),
        ("Semantic\n语义理解", Inches(3.5)),
        ("ContextManager\n上下文管理", Inches(5.0)),
        ("CoT Engine\nGPT-4推理", Inches(6.8)),
        ("Reasoning\n推理结果", Inches(8.3))
    ]

    cognitive_module_width = Inches(1.3)
    cognitive_module_height = Inches(0.7)
    cognitive_module_y = cognitive_y + Inches(0.35)

    for name, x in cognitive_modules:
        # CoT用深蓝色突出
        if "CoT" in name or "GPT" in name:
            bg_color = RGBColor(232, 234, 246)
            border_color = RGBColor(63, 81, 181)
            border_width = Pt(3)
        else:
            bg_color = RGBColor(225, 245, 254)
            border_color = RGBColor(2, 136, 209)
            border_width = Pt(1.5)

        module = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, cognitive_module_y, cognitive_module_width, cognitive_module_height
        )
        module.fill.solid()
        module.fill.fore_color.rgb = bg_color
        module.line.color.rgb = border_color
        module.line.width = border_width

        tf = module.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # ========== 规划层 ==========
    planning_y = cognitive_y + Inches(1.3)
    planning_height = Inches(1.1)

    planning_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), planning_y, Inches(9.4), planning_height
    )
    planning_box.fill.solid()
    planning_box.fill.fore_color.rgb = RGBColor(200, 230, 201)
    planning_box.line.color.rgb = RGBColor(56, 142, 60)
    planning_box.line.width = Pt(2.5)

    # 规划层标题
    planning_title = slide.shapes.add_textbox(
        Inches(0.45), planning_y + Inches(0.05), Inches(2.5), Inches(0.3)
    )
    tf = planning_title.text_frame
    tf.text = "规划层 Planning Layer (HTN分层规划)"
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 100, 50)

    # 规划层模块
    planning_modules = [
        ("TaskLevel\n任务层\n→技能序列", Inches(0.5)),
        ("SkillLevel\n技能层\n→动作序列", Inches(2.15)),
        ("ActionLevel\n动作层\n→参数化", Inches(3.8)),
        ("DynamicPlanner\n动态插入\n前置条件", Inches(5.45)),
        ("Replanning\n重规划\n失败恢复", Inches(7.1)),
        ("PlanState\n任务树", Inches(8.5))
    ]

    planning_module_width = Inches(1.4)
    planning_module_height = Inches(0.7)
    planning_module_y = planning_y + Inches(0.35)

    for name, x in planning_modules:
        module = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, planning_module_y, planning_module_width, planning_module_height
        )
        module.fill.solid()
        module.fill.fore_color.rgb = RGBColor(220, 245, 220)
        module.line.color.rgb = RGBColor(56, 142, 60)
        module.line.width = Pt(1.5)

        tf = module.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # ========== 执行层 ==========
    execution_y = planning_y + Inches(1.3)
    execution_height = Inches(1.1)

    execution_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), execution_y, Inches(9.4), execution_height
    )
    execution_box.fill.solid()
    execution_box.fill.fore_color.rgb = RGBColor(225, 190, 231)
    execution_box.line.color.rgb = RGBColor(123, 31, 162)
    execution_box.line.width = Pt(2.5)

    # 执行层标题
    execution_title = slide.shapes.add_textbox(
        Inches(0.45), execution_y + Inches(0.05), Inches(2.5), Inches(0.3)
    )
    tf = execution_title.text_frame
    tf.text = "执行层 Execution Layer (自适应执行)"
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(100, 0, 120)

    # 执行层模块
    execution_modules = [
        ("Executor\n执行器", Inches(0.5)),
        ("AdaptiveExec\n自适应\n实时监控", Inches(1.85)),
        ("Monitor\n执行监控", Inches(3.2)),
        ("FailureDetector\n失败检测", Inches(4.55)),
        ("Recovery\n恢复引擎", Inches(5.9)),
        ("Retry/Insert/\nReplan", Inches(7.25)),
        ("Observation\n结果反馈", Inches(8.5))
    ]

    execution_module_width = Inches(1.2)
    execution_module_height = Inches(0.7)
    execution_module_y = execution_y + Inches(0.35)

    for name, x in execution_modules:
        module = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, execution_module_y, execution_module_width, execution_module_height
        )
        module.fill.solid()
        module.fill.fore_color.rgb = RGBColor(243, 229, 245)
        module.line.color.rgb = RGBColor(123, 31, 162)
        module.line.width = Pt(1.5)

        tf = module.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # ========== 平台层 ==========
    platform_y = execution_y + Inches(1.3)
    platform_height = Inches(0.7)

    platform_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), platform_y, Inches(9.4), platform_height
    )
    platform_box.fill.solid()
    platform_box.fill.fore_color.rgb = RGBColor(255, 204, 188)
    platform_box.line.color.rgb = RGBColor(255, 87, 34)
    platform_box.line.width = Pt(2.5)

    # 平台层标题
    platform_title = slide.shapes.add_textbox(
        Inches(0.45), platform_y + Inches(0.05), Inches(2), Inches(0.3)
    )
    tf = platform_title.text_frame
    tf.text = "平台层 Platform Layer"
    p = tf.paragraphs[0]
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 50, 0)

    # 三个平台
    platforms = [
        ("无人机 Drone\n巡航/搜索/投送", Inches(1.5), "#4CAF50"),
        ("无人车 UGV\n巡逻/运输/抓取", Inches(4), "#2196F3"),
        ("无人船 USV\n水域搜索/检测", Inches(6.5), "#FF9800")
    ]

    platform_width = Inches(2)
    platform_module_height = Inches(0.45)
    platform_module_y = platform_y + Inches(0.15)

    for name, x, color in platforms:
        module = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, platform_module_y, platform_width, platform_module_height
        )
        module.fill.solid()
        module.fill.fore_color.rgb = RGBColor(255, 255, 255)
        module.line.color.rgb = RGBColor(255, 87, 34)
        module.line.width = Pt(2)

        tf = module.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = RGBColor(80, 80, 80)
        p.alignment = PP_ALIGN.CENTER

    # ========== 数据流说明 ==========
    dataflow_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(7.0), Inches(9.4), Inches(0.5)
    )
    tf = dataflow_box.text_frame
    tf.word_wrap = True
    tf.text = "数据流：PerceptionData → CognitiveOutput → PlanState → ExecutionResult | 关键技术：VLM视觉理解 | LLM推理决策 | WorldModel世界建模 | HTN分层规划 | 自适应执行"
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER


def create_dataflow_slide(prs):
    """创建数据流详解页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "层间数据流详解"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 四层数据结构
    layers = [
        ("感知层输出", "PerceptionData", [
            "• point_cloud: 点云数据 (xyz + intensity)",
            "• detections: 目标检测结果 [{label, bbox, confidence}]",
            "• semantic_objects: VLM语义理解 [{label, description}]",
            "• sensors_data: 传感器原始数据 {lidar, camera, imu}",
            "• timestamp: 时间戳"
        ], RGBColor(255, 235, 157)),

        ("认知层输出", "CognitiveOutput", [
            "• planning_context: PlanningContext 规划上下文",
            "  └─ robot_state: 机器人状态 {position, battery}",
            "  └─ world_objects: 世界物体列表 [{id, label, position}]",
            "  └─ spatial_relations: 空间关系 [\"杯子在桌子上\"]",
            "  └─ tracked_objects: 追踪对象 [{id, position, velocity}]",
            "• environment_changes: 环境变化列表"
        ], RGBColor(179, 229, 252)),

        ("规划层输出", "PlanState", [
            "• roots: List[PlanNode] 根节点列表",
            "• nodes: Dict[id, PlanNode] 所有节点索引",
            "• execution_history: List[Dict] 执行历史",
            "",
            "PlanNode {",
            "  • id, name, action (goto/grasp/detect)",
            "  • preconditions: 前置条件",
            "  • expected_effects: 预期效果",
            "  • parameters: 参数 {location, speed}",
            "  • status: 状态 (pending/executing/success/failed)"
        ], RGBColor(200, 230, 201)),

        ("执行层输出", "ExecutionResult", [
            "• success: bool 成功/失败",
            "• action_id: str 动作ID",
            "• result: Dict 执行结果",
            "• error: Optional[str] 错误信息",
            "",
            "失败恢复策略：",
            "• 重试: retry_count < 3",
            "• 插入: 动态插入新动作",
            "• 重规划: 触发ReplanningRules"
        ], RGBColor(225, 190, 231))
    ]

    y_start = Inches(1.0)
    layer_height = Inches(1.4)

    for i, (title, type_name, fields, color) in enumerate(layers):
        y = y_start + i * (layer_height + Inches(0.15))

        # 层框
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), y, Inches(9), layer_height
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = color
        layer_box.line.color.rgb = RGBColor(100, 100, 100)
        layer_box.line.width = Pt(1.5)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.7), y + Inches(0.08), Inches(3), Inches(0.3)
        )
        tf = title_box.text_frame
        tf.text = f"{title} → {type_name}"
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)

        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(0.7), y + Inches(0.4), Inches(8.6), Inches(0.9)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for field in fields:
            p = tf.add_paragraph()
            p.text = field
            p.font.size = Pt(10)
            p.font.family = "Courier New"
            p.font.color.rgb = RGBColor(30, 30, 30)
            p.space_before = Pt(2)
            p.space_after = Pt(1)


def create_llm_slide(prs):
    """创建大模型应用页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "大模型（VLM/LLM）应用位置"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # VLM部分
    vlm_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(0.9), Inches(4.2), Inches(2.8)
    )
    vlm_box.fill.solid()
    vlm_box.fill.fore_color.rgb = RGBColor(255, 235, 238)
    vlm_box.line.color.rgb = RGBColor(229, 57, 53)
    vlm_box.line.width = Pt(3)

    vlm_title = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.0), Inches(3.8), Inches(0.4)
    )
    tf = vlm_title.text_frame
    tf.text = "VLM 视觉语言模型"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(183, 28, 28)
    p.alignment = PP_ALIGN.CENTER

    vlm_content = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(3.8), Inches(2)
    )
    tf = vlm_content.text_frame
    tf.word_wrap = True
    tf.text = """【应用位置】感知层

【模型】
• LLaVA:7b
• MiniCPM-V
• Ollama本地部署

【功能1】场景理解
  输入: 摄像头图像RGB
  输出: 场景描述 + 物体列表 + 空间关系

【功能2】目标搜索
  输入: 图像 + 目标描述
  输出: 找到/未找到 + 位置 + 建议动作

【功能3】空间问答
  输入: 图像 + 问题("门在哪?")
  输出: 方向 + 距离 + 描述"""
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        if "【" in paragraph.text:
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(183, 28, 28)

    # LLM部分
    llm_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.3), Inches(0.9), Inches(4.2), Inches(2.8)
    )
    llm_box.fill.solid()
    llm_box.fill.fore_color.rgb = RGBColor(232, 234, 246)
    llm_box.line.color.rgb = RGBColor(63, 81, 181)
    llm_box.line.width = Pt(3)

    llm_title = slide.shapes.add_textbox(
        Inches(5.5), Inches(1.0), Inches(3.8), Inches(0.4)
    )
    tf = llm_title.text_frame
    tf.text = "LLM 大语言模型"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(48, 63, 159)
    p.alignment = PP_ALIGN.CENTER

    llm_content = slide.shapes.add_textbox(
        Inches(5.5), Inches(1.5), Inches(3.8), Inches(2)
    )
    tf = llm_content.text_frame
    tf.word_wrap = True
    tf.text = """【应用位置】认知层 CoT引擎

【模型】
• GPT-4 (OpenAI API)
• Claude (Anthropic API)
• Llama3.1 (本地Ollama)

【功能1】链式思维推理
  输入: 任务 + 环境上下文
  输出: 推理链 + 决策 + 建议
  特点: 可解释、可追溯

【功能2】任务分解
  输入: 自然语言指令
  输出: HTN任务树
  示例: "搜索灾区" → [巡航,搜索,检测,报告]

【功能3】异常处理
  输入: 失败信息 + 上下文
  输出: 原因分析 + 恢复策略"""
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(10)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        if "【" in paragraph.text:
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(48, 63, 159)

    # 数据流部分
    flow_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.0), Inches(9), Inches(3.2)
    )
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    flow_box.line.color.rgb = RGBColor(100, 100, 100)
    flow_box.line.width = Pt(2)

    flow_title = slide.shapes.add_textbox(
        Inches(0.7), Inches(4.1), Inches(8.6), Inches(0.3)
    )
    tf = flow_title.text_frame
    tf.text = "完整数据流（大模型驱动）"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(80, 80, 80)

    flow_content = slide.shapes.add_textbox(
        Inches(0.7), Inches(4.5), Inches(8.6), Inches(2.5)
    )
    tf = flow_content.text_frame
    tf.word_wrap = True
    tf.text = """步骤1: VLM场景理解
  → 摄像头图像 → VLM(LLaVA) → 场景描述 + 物体列表 + 空间关系
  → 输出示例: "前方有倒塌建筑，门位于正面中央，左侧有开放空间"

步骤2: VLM目标搜索
  → 用户指令("去建筑门口") + 图像 → VLM → 目标位置
  → 输出示例: "目标在图像中央，距离约20米，建议直行"

步骤3: 感知数据融合
  → 点云 + VLM结果 → FusionEngine → PerceptionData
  → 输出: 完整感知数据包(包含几何+语义)

步骤4: LLM推理决策
  → PerceptionData + 任务 → CoT引擎(GPT-4) → ReasoningResult
  → 输出: "步骤1:环境分析...步骤2:路径规划...决策:执行搜索任务"

步骤5: HTN任务分解
  → ReasoningResult + PlanningContext → TaskLevelPlanner → PlanState
  → 输出: HTN任务树 [goto(门口) → detect_door → open_door → enter]

步骤6: 执行与反馈
  → PlanState → Executor → 平台执行 → ObservationResult
  → 输出: 成功/失败 → 更新WorldModel → 循环"""
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(9)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        paragraph.line_spacing = 1.3


def main():
    """生成完整架构图PPT"""
    prs = Presentation()

    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 创建4页幻灯片
    create_title_slide(prs)
    create_architecture_slide(prs)
    create_dataflow_slide(prs)
    create_llm_slide(prs)

    # 保存PPT
    output_path = "/media/yangyuhui/CODES1/Brain/docs/investment_pitch/Brain_Complete_Architecture_Diagram.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页幻灯片")
    print()
    print("📋 PPT内容：")
    print("  第1页：标题页 - 展示三大技术优势")
    print("  第2页：完整架构图 - 四层详细结构（传感器→感知→认知→规划→执行→平台）")
    print("  第3页：数据流详解 - 层间接口定义（PerceptionData/CognitiveOutput/PlanState/ExecutionResult）")
    print("  第4页：大模型应用 - VLM/LLM位置和功能")


if __name__ == "__main__":
    main()
