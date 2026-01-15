#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成Brain项目技术架构PPT - 详细技术实现方案
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_architecture_diagram(slide):
    """创建完整的技术架构图"""

    # 定义颜色
    colors = {
        'perception': RGBColor(255, 224, 178),    # 浅黄
        'perception_dark': RGBColor(255, 204, 102),
        'cognitive': RGBColor(178, 235, 242),    # 浅蓝
        'cognitive_dark': RGBColor(77, 208, 225),
        'planning': RGBColor(178, 255, 178),     # 浅绿
        'planning_dark': RGBColor(77, 208, 145),
        'execution': RGBColor(230, 230, 250),    # 浅紫
        'execution_dark': RGBColor(149, 165, 211),
        'world_model': RGBColor(255, 218, 185),   # 橙色
        'platform': RGBColor(255, 183, 178)      # 红色
    }

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Brain 技术架构 - 详细实现方案"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # ========== 1. 感知层 ==========
    perception_y = Inches(1.0)

    # 感知层主框
    perception_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), perception_y, Inches(9), Inches(1.3)
    )
    perception_box.fill.solid()
    perception_box.fill.fore_color.rgb = colors['perception']
    perception_box.line.color.rgb = RGBColor(0, 0, 0)
    perception_box.line.width = Pt(2)

    # 感知层标题
    perception_title = slide.shapes.add_textbox(
        Inches(0.7), perception_y + Inches(0.05), Inches(2), Inches(0.4)
    )
    tf = perception_title.text_frame
    tf.text = "感知层 (Perception Layer)"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(102, 51, 0)

    # 感知层模块
    perception_modules = [
        ("SensorInput\n传感器输入", Inches(0.7)),
        ("PointCloudProcessor\n点云处理", Inches(2.6)),
        ("ObjectDetector\n目标检测", Inches(4.5)),
        ("FusionEngine\n融合引擎", Inches(6.4)),
        ("SituationalMap\n态势图生成", Inches(7.8))
    ]

    for name, x in perception_modules:
        module_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, perception_y + Inches(0.5), Inches(1.5), Inches(0.7)
        )
        module_box.fill.solid()
        module_box.fill.fore_color.rgb = colors['perception_dark']
        module_box.line.color.rgb = RGBColor(0, 0, 0)
        module_box.line.width = Pt(1)

        tf = module_box.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # ========== 2. 认知层 ==========
    cognitive_y = perception_y + Inches(1.5)

    # 认知层主框
    cognitive_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), cognitive_y, Inches(9), Inches(1.3)
    )
    cognitive_box.fill.solid()
    cognitive_box.fill.fore_color.rgb = colors['cognitive']
    cognitive_box.line.color.rgb = RGBColor(0, 0, 0)
    cognitive_box.line.width = Pt(2)

    # 认知层标题
    cognitive_title = slide.shapes.add_textbox(
        Inches(0.7), cognitive_y + Inches(0.05), Inches(2), Inches(0.4)
    )
    tf = cognitive_title.text_frame
    tf.text = "认知层 (Cognitive Layer)"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # 认知层模块
    cognitive_modules = [
        ("PerceptionParser\n感知解析器", Inches(0.7)),
        ("WorldModel\n世界模型", Inches(2.3)),
        ("SemanticUnderstanding\n语义理解", Inches(3.9)),
        ("ContextManager\n上下文管理", Inches(5.5)),
        ("CoTEngine\n思维链引擎", Inches(7.1))
    ]

    for name, x in cognitive_modules:
        module_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, cognitive_y + Inches(0.5), Inches(1.4), Inches(0.7)
        )
        module_box.fill.solid()
        module_box.fill.fore_color.rgb = colors['cognitive_dark']
        module_box.line.color.rgb = RGBColor(0, 0, 0)
        module_box.line.width = Pt(1)

        tf = module_box.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # ========== 3. World Model (跨层共享) ==========
    worldmodel_y = cognitive_y + Inches(1.5)

    worldmodel_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2), worldmodel_y, Inches(6), Inches(0.7)
    )
    worldmodel_box.fill.solid()
    worldmodel_box.fill.fore_color.rgb = colors['world_model']
    worldmodel_box.line.color.rgb = RGBColor(0, 0, 0)
    worldmodel_box.line.width = Pt(2)

    tf = worldmodel_box.text_frame
    tf.word_wrap = True
    tf.text = "World Model (世界模型) - 共享知识库\n几何态势 | 语义态势 | 动态态势 | 时空索引"
    p = tf.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # ========== 4. 规划层 ==========
    planning_y = worldmodel_y + Inches(1.0)

    # 规划层主框
    planning_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), planning_y, Inches(9), Inches(1.3)
    )
    planning_box.fill.solid()
    planning_box.fill.fore_color.rgb = colors['planning']
    planning_box.line.color.rgb = RGBColor(0, 0, 0)
    planning_box.line.width = Pt(2)

    # 规划层标题
    planning_title = slide.shapes.add_textbox(
        Inches(0.7), planning_y + Inches(0.05), Inches(2), Inches(0.4)
    )
    tf = planning_title.text_frame
    tf.text = "规划层 (Planning Layer)"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 0)

    # 规划层三层架构
    planning_modules = [
        ("TaskLevelPlanner\n任务层规划", Inches(0.7)),
        ("SkillLevelPlanner\n技能层规划", Inches(2.6)),
        ("ActionLevelPlanner\n动作层规划", Inches(4.5)),
        ("DynamicPlanner\n动态规划器", Inches(6.4)),
        ("ReplanningRules\n重规划规则", Inches(7.8))
    ]

    for name, x in planning_modules:
        module_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, planning_y + Inches(0.5), Inches(1.5), Inches(0.7)
        )
        module_box.fill.solid()
        module_box.fill.fore_color.rgb = colors['planning_dark']
        module_box.line.color.rgb = RGBColor(0, 0, 0)
        module_box.line.width = Pt(1)

        tf = module_box.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # ========== 5. 执行层 ==========
    execution_y = planning_y + Inches(1.5)

    # 执行层主框
    execution_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), execution_y, Inches(9), Inches(1.3)
    )
    execution_box.fill.solid()
    execution_box.fill.fore_color.rgb = colors['execution']
    execution_box.line.color.rgb = RGBColor(0, 0, 0)
    execution_box.line.width = Pt(2)

    # 执行层标题
    execution_title = slide.shapes.add_textbox(
        Inches(0.7), execution_y + Inches(0.05), Inches(2), Inches(0.4)
    )
    tf = execution_title.text_frame
    tf.text = "执行层 (Execution Layer)"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 0, 102)

    # 执行层模块
    execution_modules = [
        ("Executor\n执行器", Inches(0.7)),
        ("AdaptiveExecutor\n自适应执行器", Inches(2.3)),
        ("ExecutionMonitor\n执行监控器", Inches(3.9)),
        ("FailureDetector\n失败检测器", Inches(5.5)),
        ("StateTracker\n状态跟踪器", Inches(7.1))
    ]

    for name, x in execution_modules:
        module_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, execution_y + Inches(0.5), Inches(1.4), Inches(0.7)
        )
        module_box.fill.solid()
        module_box.fill.fore_color.rgb = colors['execution_dark']
        module_box.line.color.rgb = RGBColor(0, 0, 0)
        module_box.line.width = Pt(1)

        tf = module_box.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    # ========== 6. 平台层 ==========
    platform_y = execution_y + Inches(1.5)

    platforms = ["无人机 Drone", "无人车 UGV", "无人船 USV"]
    platform_width = Inches(2.5)
    platform_spacing = Inches(0.3)
    total_width = len(platforms) * platform_width + (len(platforms) - 1) * platform_spacing
    platform_left = Inches(5) - total_width / 2

    for i, platform in enumerate(platforms):
        platform_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            platform_left + i * (platform_width + platform_spacing),
            platform_y,
            platform_width,
            Inches(0.6)
        )
        platform_box.fill.solid()
        platform_box.fill.fore_color.rgb = colors['platform']
        platform_box.line.color.rgb = RGBColor(0, 0, 0)
        platform_box.line.width = Pt(2)

        tf = platform_box.text_frame
        tf.word_wrap = True
        tf.text = platform
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER


def create_data_flow_diagram(slide):
    """创建数据流图"""

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "数据流与核心算法"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 定义流程节点
    flows = [
        ("传感器数据\nLiDAR/相机/IMU", "0.8", "1.0"),
        ("数据预处理\n滤波/配准/增强", "0.8", "2.0"),
        ("特征提取\n点云/视觉/IMU", "0.8", "3.0"),
        ("多模态融合\n几何+语义+时序", "0.8", "4.0"),
        ("World Model\n世界模型更新", "0.8", "5.0"),
        ("CoT推理\n思维链生成", "3.0", "5.0"),
        ("HTN任务分解\n任务→技能→动作", "5.0", "5.0"),
        ("动态规划\n插入前置条件", "5.0", "3.5"),
        ("执行计划\n生成动作序列", "5.0", "2.0"),
        ("执行监控\n实时状态跟踪", "7.0", "2.0"),
        ("失败检测\n异常识别", "7.0", "3.5"),
        ("自适应恢复\n重规划/重试/回滚", "7.0", "5.0"),
    ]

    # 绘制节点
    nodes = {}
    for name, x, y in flows:
        node = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(float(x)), Inches(float(y)), Inches(1.8), Inches(0.7)
        )
        node.fill.solid()

        # 根据类型设置颜色
        if "传感器" in name or "数据" in name:
            node.fill.fore_color.rgb = RGBColor(255, 224, 178)
        elif "World Model" in name or "CoT" in name:
            node.fill.fore_color.rgb = RGBColor(178, 235, 242)
        elif "HTN" in name or "规划" in name:
            node.fill.fore_color.rgb = RGBColor(178, 255, 178)
        elif "执行" in name or "失败" in name or "恢复" in name:
            node.fill.fore_color.rgb = RGBColor(230, 230, 250)

        node.line.color.rgb = RGBColor(0, 0, 0)
        node.line.width = Pt(1)

        tf = node.text_frame
        tf.word_wrap = True
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(9)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        nodes[name] = node

    # 添加数据流箭头说明
    flow_text = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.2), Inches(9), Inches(1.0)
    )
    tf = flow_text.text_frame
    tf.word_wrap = True
    tf.text = """数据流：
传感器输入 → 预处理 → 融合 → World Model → CoT推理 → HTN规划 → 动态规划 → 执行
↑                                                                 ↓
← ← ← ← ← ← ← ← ← ← ← ← ← 失败检测 → 自适应恢复 ← ← ← ← ← ← ← ← ← ← ← ← ←"""
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(102, 0, 102)
    p.line_spacing = 1.5


def create_module_details(slide):
    """创建核心模块详解"""

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "核心模块详解"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 四列内容
    modules = [
        ("感知层", [
            "SensorInput",
            "  - 接收多传感器数据流",
            "  - 时间戳同步",
            "",
            "PointCloudProcessor",
            "  - 点云滤波与分割",
            "  - 地面平面提取",
            "",
            "ObjectDetector",
            "  - YOLO目标检测",
            "  - 3D位置估计"
        ]),
        ("认知层", [
            "WorldModel",
            "  - 几何世界建模",
            "  - 语义关系抽取",
            "",
            "SemanticUnderstanding",
            "  - VLM视觉理解",
            "  - 场景语义标注",
            "",
            "CoTEngine",
            "  - 链式推理",
            "  - 决策可解释"
        ]),
        ("规划层", [
            "TaskLevelPlanner",
            "  - 自然语言解析",
            "  - 任务树生成",
            "",
            "SkillLevelPlanner",
            "  - 技能序列分解",
            "  - 约束求解",
            "",
            "ActionLevelPlanner",
            "  - 参数化动作生成",
            "  - 前置条件检查"
        ]),
        ("执行层", [
            "AdaptiveExecutor",
            "  - 执行状态监控",
            "  - 实时调整",
            "",
            "FailureDetector",
            "  - 异常识别",
            "  - 失败分类",
            "",
            "ReplanningRules",
            "  - 重规划决策",
            "  - 恢复策略选择"
        ])
    ]

    x_positions = [Inches(0.5), Inches(3.0), Inches(5.5), Inches(8.0)]
    y_start = Inches(1.0)

    for i, (title, items) in enumerate(modules):
        # 列标题框
        header_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i], y_start, Inches(2.3), Inches(0.5)
        )

        # 颜色
        colors = [
            RGBColor(255, 204, 102),   # 感知层-黄
            RGBColor(77, 208, 225),    # 认知层-蓝
            RGBColor(77, 208, 145),    # 规划层-绿
            RGBColor(149, 165, 211)    # 执行层-紫
        ]
        header_box.fill.solid()
        header_box.fill.fore_color.rgb = colors[i]
        header_box.line.color.rgb = RGBColor(0, 0, 0)
        header_box.line.width = Pt(2)

        tf = header_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 内容
        content_box = slide.shapes.add_textbox(
            x_positions[i] + Inches(0.1), y_start + Inches(0.6), Inches(2.1), Inches(5.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            if item == "":
                p.space_after = Pt(4)
            else:
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor(0, 0, 0)
                p.space_before = Pt(2)
                p.space_after = Pt(1)


def create_tech_highlights(slide):
    """创建技术亮点"""

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "技术亮点与创新"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 技术亮点列表
    highlights = [
        ("1. World Model 驱动", [
            "多模态传感器融合（LiDAR + 视觉 + IMU）",
            "实时构建几何、语义、动态三维态势图",
            "时空索引，毫秒级查询响应"
        ]),
        ("2. Chain-of-Thought 推理", [
            "认知层显式推理链",
            "决策过程可解释、可追溯",
            "支持复杂任务分解与策略选择"
        ]),
        ("3. HTN 分层规划", [
            "三层规划器：任务层 → 技能层 → 动作层",
            "动态插入前置条件（如开门、搜索）",
            "失败自动恢复（重试/插入/重规划）"
        ]),
        ("4. 自适应执行", [
            "实时监控执行状态",
            "异常自动检测与分类",
            "多级恢复策略（回滚/重规划/求助）"
        ]),
        ("5. 平台抽象层", [
            "统一Capability接口",
            "一次开发，三平台复用（无人机/车/船）",
            "代码复用率90%+"
        ])
    ]

    y_start = Inches(1.0)
    for title, items in highlights:
        # 标题框
        title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), y_start, Inches(9), Inches(0.4)
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = RGBColor(0, 102, 204)
        title_box.line.color.rgb = RGBColor(0, 0, 0)

        tf = title_box.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(0.7), y_start + Inches(0.5), Inches(8.6), Inches(0.7)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0, 0, 0)
            p.space_before = Pt(2)

        y_start += Inches(1.3)


def main():
    """生成技术架构PPT"""
    prs = Presentation()

    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 创建幻灯片
    create_architecture_diagram(prs.slides.add_slide(prs.slide_layouts[6]))
    create_data_flow_diagram(prs.slides.add_slide(prs.slide_layouts[6]))
    create_module_details(prs.slides.add_slide(prs.slide_layouts[6]))
    create_tech_highlights(prs.slides.add_slide(prs.slide_layouts[6]))

    # 保存PPT
    output_path = "/media/yangyuhui/CODES1/Brain/docs/investment_pitch/Brain_Technical_Architecture.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页幻灯片")


if __name__ == "__main__":
    main()
