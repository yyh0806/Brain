#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成Brain项目技术实现详细方案PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_slide_1_data_flow(prs):
    """第1页：数据流详解"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "各层接口与数据流"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 感知层输出
    y = Inches(1.2)

    layer_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(9), Inches(1.3)
    )
    layer_box.fill.solid()
    layer_box.fill.fore_color.rgb = RGBColor(255, 224, 178)
    layer_box.line.color.rgb = RGBColor(0, 0, 0)
    layer_box.line.width = Pt(2)

    tf = layer_box.text_frame
    tf.word_wrap = True
    tf.text = "感知层输出 → 认知层"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.35), Inches(8.6), Inches(0.8)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = """PerceptionData {
  • point_cloud: 点云数据 (xyz + intensity)
  • detections: 目标检测结果 [{label, bbox, confidence, position_3d}]
  • semantic_objects: VLM语义理解结果 [{label, description, bbox, attributes}]
  • sensors_data: 传感器原始数据 {lidar, camera, imu, gps}
  • timestamp: 时间戳
}

关键接口：
process_perception(perception_data) -> CognitiveOutput"""
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0, 0, 0)

    # 认知层输出
    y = Inches(2.7)

    layer_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(9), Inches(1.3)
    )
    layer_box.fill.solid()
    layer_box.fill.fore_color.rgb = RGBColor(178, 235, 242)
    layer_box.line.color.rgb = RGBColor(0, 0, 0)
    layer_box.line.width = Pt(2)

    tf = layer_box.text_frame
    tf.word_wrap = True
    tf.text = "认知层输出 → 规划层"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.35), Inches(8.6), Inches(0.8)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = """CognitiveOutput {
  • planning_context: PlanningContext
      - robot_state: 机器人状态 {position, battery, capabilities}
      - world_objects: 世界物体列表 [{id, label, position, attributes}]
      - spatial_relations: 空间关系 ["杯子在桌子上"]
      - tracked_objects: 追踪对象 [{id, position, velocity, history}]
  • environment_changes: 环境变化列表
  • timestamp: 时间戳
}

关键接口：
get_planning_context() -> PlanningContext
reason(query, context, mode) -> ReasoningResult"""
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0, 0, 0)

    # 规划层输出
    y = Inches(4.2)

    layer_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(9), Inches(1.3)
    )
    layer_box.fill.solid()
    layer_box.fill.fore_color.rgb = RGBColor(178, 255, 178)
    layer_box.line.color.rgb = RGBColor(0, 0, 0)
    layer_box.line.width = Pt(2)

    tf = layer_box.text_frame
    tf.word_wrap = True
    tf.text = "规划层输出 → 执行层"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.35), Inches(8.6), Inches(0.8)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = """PlanState {
  • roots: List[PlanNode]  # 根节点列表
  • nodes: Dict[id, PlanNode]  # 所有节点索引
  • execution_history: List[Dict]  # 执行历史
}

PlanNode {
  • id: 节点ID
  • name: 节点名称
  • action: 动作类型 (goto/grasp/detect/...)
  • preconditions: List[str]  # 前置条件
  • expected_effects: List[str]  # 预期效果
  • parameters: Dict  # 参数
  • status: NodeStatus (pending/executing/success/failed)
  • children: List[PlanNode]  # 子节点
}

关键接口：
get_plan(command) -> PlanState
plan_and_execute(command, robot_interface) -> result"""
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0, 0, 0)

    # 执行层输出
    y = Inches(5.7)

    layer_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(9), Inches(1.0)
    )
    layer_box.fill.solid()
    layer_box.fill.fore_color.rgb = RGBColor(230, 230, 250)
    layer_box.line.color.rgb = RGBColor(0, 0, 0)
    layer_box.line.width = Pt(2)

    tf = layer_box.text_frame
    tf.word_wrap = True
    tf.text = "执行层 → 平台"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.2), Inches(8.6), Inches(0.6)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = """ExecutionResult {
  • success: bool
  • action_id: str
  • result: Dict  # 执行结果
  • error: Optional[str]
}

关键接口：
execute_plan(plan_state, robot_interface) -> ExecutionResult"""
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0, 0, 0)


def create_slide_2_llm_applications(prs):
    """第2页：大模型应用"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "大模型（LLM/VLM）应用位置"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # VLM应用
    y = Inches(1.0)

    vlm_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(4.2), Inches(2.8)
    )
    vlm_box.fill.solid()
    vlm_box.fill.fore_color.rgb = RGBColor(255, 235, 205)
    vlm_box.line.color.rgb = RGBColor(0, 0, 0)
    vlm_box.line.width = Pt(2)

    tf = vlm_box.text_frame
    tf.word_wrap = True
    tf.text = "1. VLM (视觉语言模型)"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(102, 0, 0)
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.5), Inches(3.8), Inches(2.2)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    items = [
        ("应用位置", "感知层"),
        ("", ""),
        ("模型", "LLaVA, MiniCPM-V"),
        ("(Ollama本地部署)", ""),
        ("", ""),
        ("功能1", "场景理解"),
        ("输入", "摄像头图像RGB"),
        ("输出", "场景描述 + 物体列表 + 空间关系"),
        ("", ""),
        ("功能2", "目标搜索"),
        ("输入", "图像 + 目标描述"),
        ("输出", "找到/未找到 + 位置 + 建议动作"),
        ("", ""),
        ("功能3", "空间问答"),
        ("输入", "图像 + 问题(门在哪?)"),
        ("输出", "方向 + 距离 + 描述")
    ]

    for title, content in items:
        if title == "":
            p = tf.add_paragraph()
            p.text = ""
            p.space_after = Pt(2)
        else:
            p = tf.add_paragraph()
            p.text = f"{title}: {content}"
            if title == "应用位置":
                p.font.bold = True
                p.font.color.rgb = RGBColor(204, 0, 0)
            elif title == "功能1" or title == "功能2" or title == "功能3":
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 51, 0)
            else:
                p.font.size = Pt(10)
            p.space_before = Pt(2)
            p.space_after = Pt(2)

    # LLM应用
    llm_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.3), y, Inches(4.2), Inches(2.8)
    )
    llm_box.fill.solid()
    llm_box.fill.fore_color.rgb = RGBColor(205, 235, 255)
    llm_box.line.color.rgb = RGBColor(0, 0, 0)
    llm_box.line.width = Pt(2)

    tf = llm_box.text_frame
    tf.word_wrap = True
    tf.text = "2. LLM (大语言模型)"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 102)
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(5.5), y + Inches(0.5), Inches(3.8), Inches(2.2)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    items = [
        ("应用位置", "认知层 CoT引擎"),
        ("", ""),
        ("模型", "GPT-4, Claude, 本地Llama"),
        ("(API调用或本地部署)", ""),
        ("", ""),
        ("功能1", "链式思维推理"),
        ("输入", "任务 + 环境上下文"),
        ("输出", "推理链 + 决策 + 建议"),
        ("", ""),
        ("功能2", "任务分解"),
        ("输入", "自然语言指令"),
        ("输出", "HTN任务树"),
        ("", ""),
        ("功能3", "异常处理"),
        ("输入", "失败信息 + 上下文"),
        ("输出", "原因分析 + 恢复策略"),
        ("", ""),
        ("核心价值", "可解释性、可追溯")
    ]

    for title, content in items:
        if title == "":
            p = tf.add_paragraph()
            p.text = ""
            p.space_after = Pt(2)
        else:
            p = tf.add_paragraph()
            p.text = f"{title}: {content}"
            if title == "应用位置":
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 204)
            elif title.startswith("功能"):
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 51, 102)
            elif title == "核心价值":
                p.font.bold = True
                p.font.color.rgb = RGBColor(204, 0, 0)
            else:
                p.font.size = Pt(10)
            p.space_before = Pt(2)
            p.space_after = Pt(2)

    # 数据流
    y = Inches(4.2)

    flow_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(9), Inches(2.8)
    )
    flow_box.fill.solid()
    flow_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    flow_box.line.color.rgb = RGBColor(0, 0, 0)
    flow_box.line.width = Pt(2)

    tf = flow_box.text_frame
    tf.word_wrap = True
    tf.text = "3. 完整数据流（大模型驱动）"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.4), Inches(8.6), Inches(2.2)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = """步骤1: VLM场景理解
  摄像头图像 → VLM(LLaVA) → 场景描述 + 物体列表
  输出: "前方有建筑、门、道路，门位于建筑正面中央"

步骤2: VLM目标搜索
  用户指令:"去建筑门口" + 图像 → VLM → 目标位置
  输出: "目标在图像中央，距离约20米，建议直行"

步骤3: 感知数据融合
  点云 + VLM结果 → 感知融合模块 → PerceptionData
  输出: 完整的感知数据包(包含几何+语义)

步骤4: LLM推理决策
  PerceptionData + 任务 → CoT引擎(LLM) → ReasoningResult
  输出: "步骤1:环境分析...步骤2:路径规划...决策:执行任务"

步骤5: HTN任务分解
  ReasoningResult + PlanningContext → TaskLevelPlanner → PlanState
  输出: HTN任务树 [goto(门口) → detect_door → open_door → enter]

步骤6: 执行与反馈
  PlanState → Executor → 平台执行 → ObservationResult
  输出: 成功/失败 → 更新WorldModel → 循环"""
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.line_spacing = 1.3


def create_slide_3_implementation(prs):
    """第3页：具体实现方案"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "具体实现方案：如何做？"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 三列实现步骤
    implementations = [
        ("阶段1: VLM感知", [
            "部署 Ollama + LLaVA",
            "  docker run -d -v ollama:/root/.ollama -p 11434:11434 ollama/ollama",
            "  ollama pull llava:7b",
            "",
            "场景理解",
            "  vlm.describe_scene(image_rgb) -> SceneDescription",
            "  返回: {summary, objects[], spatial_relations[], hints[]}",
            "",
            "目标搜索",
            "  vlm.find_target(image, \"红色杯子\") -> TargetSearchResult",
            "  返回: {found, position, confidence, action}"
        ]),
        ("阶段2: LLM推理", [
            "配置 LLM 接口",
            "  - 本地: Ollama + Llama3.1",
            "  - 云端: OpenAI GPT-4 API",
            "",
            "CoT推理",
            "  cot.reason(\"去厨房拿杯子\", context, PLANNING)",
            "  返回: ReasoningResult{chain[], decision, suggestion}",
            "",
            "任务分解",
            "  task_planner.parse_command(\"去厨房拿杯水\")",
            "  返回: TaskInfo{skills:[\"navigate\", \"search\", \"grasp\"]}"
        ]),
        ("阶段3: HTN规划", [
            "三层规划器",
            "  TaskLevel: 自然语言 → 技能序列",
            "  SkillLevel: 技能 → 动作序列",
            "  ActionLevel: 动作 → 参数化操作",
            "",
            "动态规划",
            "  dynamic_planner.check_and_insert_preconditions(node)",
            "  检测到门关闭 → 自动插入 open_door 动作",
            "",
            "失败恢复",
            "  replanning_rules.should_replan(failed_node, counts)",
            "  超过3次插入 → 触发重规划"
        ])
    ]

    x_positions = [Inches(0.5), Inches(3.5), Inches(6.5)]
    y_start = Inches(1.0)

    for i, (title, items) in enumerate(implementations):
        # 列框
        col_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i], y_start, Inches(2.8), Inches(6.2)
        )

        colors = [
            RGBColor(255, 235, 205),  # 橙色
            RGBColor(205, 235, 255),  # 蓝色
            RGBColor(205, 255, 205)   # 绿色
        ]
        col_box.fill.solid()
        col_box.fill.fore_color.rgb = colors[i]
        col_box.line.color.rgb = RGBColor(0, 0, 0)
        col_box.line.width = Pt(2)

        # 标题
        title_box = slide.shapes.add_textbox(
            x_positions[i] + Inches(0.1), y_start + Inches(0.1), Inches(2.6), Inches(0.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 内容
        content_box = slide.shapes.add_textbox(
            x_positions[i] + Inches(0.1), y_start + Inches(0.7), Inches(2.6), Inches(5.4)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            if item.startswith("部署") or item.startswith("配置") or item.startswith("三层"):
                p.font.bold = True
                p.font.color.rgb = RGBColor(102, 0, 0)
            elif "阶段" not in item:
                p.font.size = Pt(9)
                p.font.family = "Courier New"
            else:
                p.font.size = Pt(11)
            p.space_before = Pt(2)
            p.space_after = Pt(1)


def main():
    """生成技术实现PPT"""
    prs = Presentation()

    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 创建3页幻灯片
    create_slide_1_data_flow(prs)
    create_slide_2_llm_applications(prs)
    create_slide_3_implementation(prs)

    # 保存PPT
    output_path = "/media/yangyuhui/CODES1/Brain/docs/investment_pitch/Brain_Implementation_Detail.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页幻灯片")


if __name__ == "__main__":
    main()
