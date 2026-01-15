#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Brain系统任务执行流程详解 - 搜救灾区被困人员
展示各层如何协作完成复杂任务
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


def create_slide_task_breakdown(prs):
    """第1页：任务完整流程"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "案例：搜救灾区被困人员"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "完整任务：用户指令 \"搜索灾区，发现被困人员后立即报告\""
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    # 四层架构的职责
    y = Inches(1.4)
    layer_height = Inches(1.3)

    layers = [
        ("感知层: 看到什么？", RGBColor(255, 224, 178), [
            "输入: 传感器数据",
            "  • 相机RGB图像",
            "  • 激光雷达点云",
            "  • IMU姿态",
            "  • GPS位置",
            "",
            "VLM场景理解:",
            "  「这是一片倒塌的建筑区域，有废墟、",
            "   碎石、开放空间，部分建筑结构完整」",
            "",
            "目标检测:",
            "  • 检测到3个可能的目标",
            "  • 热成像发现1个热源",
            "",
            "输出: PerceptionData",
            "  {场景描述, 物体列表, 空间关系}"
        ]),
        ("认知层: 理解什么？", RGBColor(178, 235, 242), [
            "输入: PerceptionData",
            "",
            "WorldModel构建:",
            "  • 几何世界: 建筑物3D模型、废墟分布",
            "  • 语义世界: 「倒塌建筑」「可能的避难所」",
            "  • 动态世界: 跟踪3个潜在目标的位置",
            "",
            "CoT推理:",
            "  步骤1: 灾区环境复杂，优先搜索开放空间",
            "  步骤2: 热源可能是被困人员，优先级最高",
            "  步骤3: 规划搜索路径：中心→外围→返航",
            "",
            "输出: PlanningContext",
            "  {机器人状态, 世界物体, 空间关系, 推理链}"
        ]),
        ("规划层: 怎么做？", RGBColor(178, 255, 178), [
            "输入: PlanningContext + \"搜索灾区，发现被困人员\"",
            "",
            "任务层规划:",
            "  自然语言 → 任务树",
            "  [搜索灾区] → [搜索开放区域] → [检测热源] → [确认人员]",
            "",
            "技能层规划:",
            "  [搜索开放区域] → skills: [起飞, 巡航, 搜索]",
            "  [检测热源] → skills: [热成像, 识别, 定位]",
            "",
            "动作层规划:",
            "  [起飞] → action: takeoff(params={height: 50m})",
            "  [巡航] → action: goto(params={path: [...]})",
            "  [搜索] → action: scan(params={mode: thermal})",
            "",
            "输出: PlanState (HTN任务树)"
        ]),
        ("执行层: 如何执行？", RGBColor(230, 230, 250), [
            "输入: PlanState",
            "",
            "执行过程:",
            "  1. takeoff(50m) → 执行中 → 成功 → 标记SUCCESS",
            "  2. goto(search_point) → 执行中 → 成功 → 标记SUCCESS",
            "  3. scan(thermal) → 执行中 → 发现热源 → 标记SUCCESS",
            "  4. detect(person) → 执行中 → 确认 → 标记SUCCESS",
            "",
            "异常处理:",
            "  • 遇到强风 → goto失败 → 重试1次 → 调整高度",
            "  • 热源消失 → scan失败 → 插入search → 重新扫描",
            "",
            "输出: ExecutionResult",
            "  {success: true, location: (x,y,z), action: \"报告\"}"
        ])
    ]

    x_positions = [Inches(0.5), Inches(2.8), Inches(5.1), Inches(7.4)]

    for i, (title, color, content_lines) in enumerate(layers):
        # 层框
        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[i], y, Inches(2.2), layer_height
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = color
        layer_box.line.color.rgb = RGBColor(0, 0, 0)
        layer_box.line.width = Pt(2)

        # 标题
        title_box_inner = slide.shapes.add_textbox(
            x_positions[i] + Inches(0.1), y + Inches(0.05), Inches(2), Inches(0.4)
        )
        tf = title_box_inner.text_frame
        tf.word_wrap = True
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 内容
        content_box = slide.shapes.add_textbox(
            x_positions[i] + Inches(0.1), y + Inches(0.45), Inches(2), Inches(0.8)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for line in content_lines:
            p = tf.add_paragraph()
            p.text = line
            if line.startswith("输入:") or line.startswith("输出:"):
                p.font.bold = True
                p.font.color.rgb = RGBColor(0, 0, 102)
            elif "VLM" in line or "CoT" in line or "HTN" in line:
                p.font.bold = True
                p.font.color.rgb = RGBColor(204, 0, 0)
            else:
                p.font.size = Pt(8)
            p.space_before = Pt(1)
            p.space_after = Pt(1)

    # 数据流箭头
    for i in range(3):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x_positions[i] + Inches(2.2), y + Inches(0.55), Inches(0.15), Inches(0.2)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
        arrow.line.fill.background()

    # 底部总结
    summary_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.9), Inches(9), Inches(0.6)
    )
    tf = summary_box.text_frame
    tf.word_wrap = True
    tf.text = """数据流: 传感器 → VLM理解 → WorldModel → CoT推理 → HTN规划 → 自适应执行 → 平台控制
关键: 每层只做自己的事，上层不干预下层的具体实现"""
    p = tf.paragraphs[0]
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(102, 0, 102)
    p.alignment = PP_ALIGN.CENTER


def create_slide_worldmodel(prs):
    """第2页：World Model 智能理解"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "第一层：World Model 智能理解"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "任务：将感知数据转化为可理解的世界模型"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    # 左侧：输入输出
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.4), Inches(4), Inches(5.5)
    )
    tf = left_box.text_frame
    tf.word_wrap = True

    # 输入
    p = tf.paragraphs[0]
    p.text = "📥 输入：PerceptionData"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(204, 0, 0)
    p.space_after = Pt(8)

    input_data = [
        "原始数据:",
        "  point_cloud: 点云数据 (100万+点)",
        "  detections: YOLO检测结果 [",
        "    {label: \"person\", bbox: [...], confidence: 0.85},",
        "    {label: \"rubble\", bbox: [...], confidence: 0.92},",
        "    {label: \"building\", bbox: [...], confidence: 0.78}",
        "  ]",
        "  semantic_objects: VLM理解结果 [",
        "    {label: \"倒塌建筑\", description: \"部分结构受损\",",
        "     bbox: [...], confidence: 0.7},",
        "    {label: \"开放空间\", description: \"可能的避难所\",",
        "     position: \"中央区域\", confidence: 0.8}",
        "  ]"
    ]

    for line in input_data:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.family = "Courier New"
        if line.startswith("  ") and ":" in line and not line.startswith("    "):
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 0)
        p.space_before = Pt(2)
        p.space_after = Pt(2)

    # World Model处理
    y = Inches(3.0)

    wm_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(4), Inches(1.5)
    )
    wm_box.fill.solid()
    wm_box.fill.fore_color.rgb = RGBColor(255, 248, 220)
    wm_box.line.color.rgb = RGBColor(255, 152, 0)
    wm_box.line.width = Pt(3)

    tf = wm_box.text_frame
    tf.word_wrap = True
    tf.text = "⚙️ WorldModel 处理"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.4), Inches(3.6), Inches(1.0)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = """1. 更新几何世界：点云 → 3D建筑模型
2. 更新语义世界：「倒塌建筑」「废墟」
3. 更新动态世界：跟踪3个潜在目标
4. 建立空间关系：「废墟在建筑左侧」
5. 生成PlanningContext"""
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.line_spacing = 1.4

    # 输出
    y = Inches(4.7)

    p = left_box.text_frame.add_paragraph()
    p.text = "📤 输出：PlanningContext"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 0)
    p.space_before = Pt(12)
    p.space_after = Pt(8)

    output_data = [
        "机器人状态:",
        "  robot_state: {",
        "    position: {x: 100, y: 200, z: 50},",
        "    battery: 85%,",
        "    capabilities: [\"fly\", \"thermal_camera\"]",
        "  }",
        "",
        "世界物体:",
        "  world_objects: [",
        "    {id: \"building_1\", label: \"倒塌建筑\",",
        "     position: {x: 150, y: 250}, status: \"damaged\"},",
        "    {id: \"target_1\", label: \"热源\",",
        "     position: {x: 120, y: 220}, confidence: 0.9}",
        "  ]",
        "",
        "空间关系:",
        "  spatial_relations: [",
        "    \"热源在倒塌建筑内部\",",
        "    \"废墟阻挡了直接路径\"",
        "  ]",
        "",
        "追踪对象:",
        "  tracked_objects: [",
        "    {id: \"target_1\", position: {...}, velocity: {...},",
        "     history: [...], last_seen: \"2秒前\"}",
        "  ]"
    ]

    for line in output_data:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(9)
        p.font.family = "Courier New"
        if line.startswith("  ") and ":" in line and not line.startswith("    {"):
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 102)
        p.space_before = Pt(2)
        p.space_after = Pt(2)

    # 右侧：核心能力
    right_box = slide.shapes.add_textbox(
        Inches(5.0), Inches(1.4), Inches(4.5), Inches(5.5)
    )
    tf = right_box.text_frame
    tf.word_wrap = True

    capabilities = [
        ("🎯 核心能力1: 几何世界建模", [
            "输入: 点云数据 (100万+点)",
            "处理: 地面分割 → 障碍物提取 → 3D重建",
            "输出: 3D建筑模型、可通行区域地图",
            "",
            "价值: 规划层知道「哪里能飞」「哪里有障碍」"
        ]),
        ("🧠 核心能力2: 语义世界理解", [
            "输入: VLM场景描述 + YOLO检测结果",
            "处理: 语义标注 → 关系抽取 → 场景分类",
            "输出: 「倒塌建筑」「可能的避难所」「废墟」",
            "",
            "价值: 规划层知道「这是什么」「有什么意义」"
        ]),
        ("📍 核心能力3: 动态世界跟踪", [
            "输入: 连续帧的检测结果",
            "处理: 目标关联 → 位置预测 → 状态更新",
            "输出: 每个目标的历史轨迹、速度、方向",
            "",
            "价值: 规划层知道「目标在哪」「往哪移动」"
        ]),
        ("🔗 核心能力4: 空间关系推理", [
            "输入: 物体位置 + 几何模型",
            "处理: 拓扑关系 → 遮挡关系 → 可达性分析",
            "输出: 「A在B内部」「C被D阻挡」",
            "",
            "价值: 规划层知道「能做什么」「不能做什么」"
        ])
    ]

    y_start = Inches(1.4)
    for title, items in capabilities:
        # 能力框
        cap_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(5.2), y_start, Inches(4.1), Inches(1.2)
        )
        cap_box.fill.solid()
        cap_box.fill.fore_color.rgb = RGBColor(240, 248, 255)
        cap_box.line.color.rgb = RGBColor(0, 102, 204)
        cap_box.line.width = Pt(1)

        tf = cap_box.text_frame
        tf.word_wrap = True
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(5.3), y_start + Inches(0.2), Inches(3.9), Inches(1.0)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(9)
            p.space_before = Pt(2)
            p.space_after = Pt(1)

        y_start += Inches(1.25)


def create_slide_htn(prs):
    """第3页：HTN分层规划 + 动态推理"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "第二层：HTN分层规划 + 动态推理"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "任务：将用户指令转化为可执行的HTN任务树"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    # 三层规划
    y = Inches(1.4)

    # 任务层
    task_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(9), Inches(1.1)
    )
    task_box.fill.solid()
    task_box.fill.fore_color.rgb = RGBColor(179, 229, 252)
    task_box.line.color.rgb = RGBColor(0, 0, 102)
    task_box.line.width = Pt(2)

    tf = task_box.text_frame
    tf.word_wrap = True
    tf.text = "任务层规划 (TaskLevelPlanner): 自然语言 → 任务树"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.15), Inches(8.6), Inches(0.85)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = """输入: \"搜索灾区，发现被困人员后立即报告\"
输出: TaskNode {
  id: "search_rescue",
  name: "搜索救援",
  type: "compound",
  children: [
    TaskNode("搜索灾区", "search_area"),
    TaskNode("检测人员", "detect_person"),
    TaskNode("报告位置", "report_location")
  ]
}
关键: 将模糊的自然语言拆解为明确的子任务"""
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.family = "Courier New"

    # 技能层
    y = Inches(2.7)

    skill_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(9), Inches(1.3)
    )
    skill_box.fill.solid()
    skill_box.fill.fore_color.rgb = RGBColor(167, 230, 219)
    skill_box.line.color.rgb = RGBColor(0, 51, 51)
    skill_box.line.width = Pt(2)

    tf = skill_box.text_frame
    tf.word_wrap = True
    tf.text = "技能层规划 (SkillLevelPlanner): 任务 → 技能序列"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.15), Inches(8.6), Inches(1.05)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = """输入: TaskNode("搜索灾区")
输出: SkillNode {
  id: "search_area_1",
  name: "区域搜索",
  skill: "aerial_search",
  children: [
    SkillNode("起飞", "takeoff", skill="flight_control"),
    SkillNode("巡航", "cruise", skill="flight_control"),
    SkillNode("搜索", "search", skill="visual_search"),
    SkillNode("检测", "detect", skill="thermal_detection")
  ]
}
关键: 选择可用的技能（Capability），检查平台支持"""
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.family = "Courier New"

    # 动作层
    y = Inches(4.3)

    action_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(9), Inches(1.3)
    )
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = RGBColor(165, 214, 167)
    action_box.line.color.rgb = RGBColor(0, 51, 0)
    action_box.line.width = Pt(2)

    tf = action_box.text_frame
    tf.word_wrap = True
    tf.text = "动作层规划 (ActionLevelPlanner): 技能 → 参数化动作"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.15), Inches(8.6), Inches(1.05)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = """输入: SkillNode("起飞")
输出: ActionNode {
  id: "takeoff_1",
  action: "takeoff",
  parameters: {height: 50, speed: 5, mode: "vertical"},
  preconditions: ["battery > 20%", "gps_ready"],
  expected_effects: ["altitude == 50m"]
}
关键: 生成具体参数，检查前置条件"""
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.family = "Courier New"

    # 动态规划
    y = Inches(5.9)

    dynamic_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), y, Inches(9), Inches(0.9)
    )
    dynamic_box.fill.solid()
    dynamic_box.fill.fore_color.rgb = RGBColor(255, 235, 205)
    dynamic_box.line.color.rgb = RGBColor(204, 51, 0)
    dynamic_box.line.width = Pt(2)

    tf = dynamic_box.text_frame
    tf.word_wrap = True
    tf.text = "动态规划 (DynamicPlanner): 运行时插入必要操作"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    content_box = slide.shapes.add_textbox(
        Inches(0.7), y + Inches(0.1), Inches(8.6), Inches(0.7)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.text = """场景1: 起飞前检测电池不足 → 动态插入 charge_action
场景2: 检测到门/障碍 → 动态插入 open_door / remove_obstacle
场景3: 热源消失 → 动态调整搜索路径 → 重新规划
关键: 根据环境变化实时调整计划"""
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.family = "Courier New"


def create_slide_execution(prs):
    """第4页：自适应执行引擎"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = "第三层：自适应执行引擎"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.8), Inches(9), Inches(0.4)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "任务：执行HTN任务树，处理异常，自动恢复"
    p = subtitle_frame.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    # 左侧：执行流程
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.4), Inches(5), Inches(5.5)
    )
    tf = left_box.text_frame
    tf.word_wrap = True

    # 执行流程标题
    p = tf.paragraphs[0]
    p.text = "📋 执行流程（5个节点示例）"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 0)
    p.space_after = Pt(8)

    executions = [
        ("✅ 节点1: takeoff(50m)", [
            "状态: PENDING → EXECUTING → SUCCESS",
            "操作: 无人机垂直起飞到50米",
            "耗时: 8秒",
            "结果: 成功到达目标高度"
        ]),
        ("✅ 节点2: goto(search_point)", [
            "状态: PENDING → EXECUTING → SUCCESS",
            "操作: 沿规划路径飞行到搜索点",
            "耗时: 15秒",
            "结果: 成功到达"
        ]),
        ("⚠️  节点3: scan(thermal)", [
            "状态: PENDING → EXECUTING → FAILED",
            "操作: 热成像扫描",
            "错误: 未发现明显热源",
            "恢复: DynamicPlanner.insert(search)",
            "      → 重新扫描，扩大范围"
        ]),
        ("✅ 节点4: detect(person)", [
            "状态: PENDING → EXECUTING → SUCCESS",
            "操作: 检测人员（热像+视觉）",
            "耗时: 5秒",
            "结果: 发现被困人员！",
            "位置: (120, 220, 10), 置信度: 0.92"
        ]),
        ("✅ 节点5: report(location)", [
            "状态: PENDING → EXECUTING → SUCCESS",
            "操作: 返航并报告位置",
            "耗时: 12秒",
            "结果: 任务完成"
        ])
    ]

    y_start = Inches(2.0)
    for i, (title, details) in enumerate(executions):
        # 节点框
        node_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.7), y_start, Inches(4.6), Inches(0.7)
        )
        node_box.fill.solid()
        if "✅" in title:
            node_box.fill.fore_color.rgb = RGBColor(200, 255, 200)
        elif "⚠️" in title:
            node_box.fill.fore_color.rgb = RGBColor(255, 220, 150)
        node_box.line.color.rgb = RGBColor(0, 0, 0)
        node_box.line.width = Pt(1)

        tf = node_box.text_frame
        tf.word_wrap = True
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True

        # 详情
        detail_box = slide.shapes.add_textbox(
            Inches(0.7), y_start + Inches(0.75), Inches(4.6), Inches(0.7)
        )
        tf = detail_box.text_frame
        tf.word_wrap = True

        for detail in details:
            p = tf.add_paragraph()
            p.text = detail
            p.font.size = Pt(9)
            p.space_before = Pt(1)
            p.space_after = Pt(1)

        y_start += Inches(1.5)

    # 右侧：核心能力
    right_box = slide.shapes.add_textbox(
        Inches(5.8), Inches(1.4), Inches(3.7), Inches(5.5)
    )
    tf = right_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "🔧 核心能力"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(102, 0, 102)
    p.space_after = Pt(10)

    capabilities = [
        ("1. 实时监控", [
            "监听每个节点的执行状态",
            "检测异常（超时、失败、偏离）",
            "更新WorldModel中的物体状态"
        ]),
        ("2. 失败检测", [
            "FailureType分类:",
            "  • PRECONDITION_FAILED  # 前置条件不满足",
            "  • EXECUTION_FAILED     # 执行失败",
            "  • WORLD_STATE_CHANGED  # 环境突变",
            "  • TIMEOUT              # 超时"
        ]),
        ("3. 自动恢复", [
            "策略1: 重试 (retry)",
            "  → 适用于临时故障",
            "",
            "策略2: 动态插入 (insert)",
            "  → 检测到门关闭，插入open_door",
            "  → 检测到物体丢失，插入search",
            "",
            "策略3: 重规划 (replan)",
            "  → 超过3次插入/重试失败",
            "  → 环境发生重大变化",
            "  → 目标不可达"
        ]),
        ("4. 状态管理", [
            "维护PlanNode的状态机:",
            "  PENDING → EXECUTING → SUCCESS/FAILED",
            "",
            "记录执行历史:",
            "  timestamp, node_id, status, result",
            "",
            "生成执行统计:",
            "  total, successful, failed, success_rate"
        ])
    ]

    y_start = Inches(2.0)
    for title, items in capabilities:
        # 能力框
        cap_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.0), y_start, Inches(3.3), Inches(1.0)
        )
        cap_box.fill.solid()
        cap_box.fill.fore_color.rgb = RGBColor(245, 245, 220)
        cap_box.line.color.rgb = RGBColor(102, 102, 0)
        cap_box.line.width = Pt(1)

        tf = cap_box.text_frame
        tf.word_wrap = True
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(12)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 内容
        content_box = slide.shapes.add_textbox(
            Inches(6.1), y_start + Inches(0.2), Inches(3.1), Inches(0.8)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(9)
            p.space_before = Pt(2)
            p.space_after = Pt(1)

        y_start += Inches(1.3)


def main():
    """生成PPT"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_slide_task_breakdown(prs)
    create_slide_worldmodel(prs)
    create_slide_htn(prs)
    create_slide_execution(prs)

    output_path = "/media/yangyuhui/CODES1/Brain/docs/investment_pitch/Brain_Task_Execution_Detail.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页幻灯片")


if __name__ == "__main__":
    main()
