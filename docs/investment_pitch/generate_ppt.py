#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成Brain项目投资路演PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_title_slide(prs):
    """创建封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)  # 深色背景
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Brain"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)  # 科技蓝
    title_para.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.5), Inches(9), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "通用小微特机器人智能操作系统"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_para.alignment = PP_ALIGN.CENTER

    # 定位语
    tagline_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5), Inches(9), Inches(0.8)
    )
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "做机器人领域的 Android / iOS"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(24)
    tagline_para.font.color.rgb = RGBColor(255, 102, 0)  # 活力橙
    tagline_para.alignment = PP_ALIGN.CENTER

    # 市场规模
    market_box = slide.shapes.add_textbox(
        Inches(2), Inches(6.5), Inches(6), Inches(1)
    )
    market_frame = market_box.text_frame
    market_frame.word_wrap = True
    market_frame.text = "🎯 2025年中国小微特机器人市场规模：3000亿元"
    market_para = market_frame.paragraphs[0]
    market_para.font.size = Pt(20)
    market_para.font.color.rgb = RGBColor(200, 200, 200)
    market_para.alignment = PP_ALIGN.CENTER


def create_why_slide(prs):
    """创建为什么做这个项目页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "为什么做这个项目？"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 左侧：市场痛点
    left_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(4.2), Inches(5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    p.text = "市场痛点"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 102, 0)
    p.space_after = Pt(20)

    pains = [
        "❌ 开发效率低\n   每个平台从零开发\n   周期6-12个月",
        "❌ 智能化不足\n   只能执行预设任务\n   缺乏环境理解",
        "❌ 协同困难\n   异构平台无法统一调度\n   多平台协作依赖人工",
        "❌ 缺乏通用OS\n   无成熟通用系统\n   代码无法复用"
    ]

    for pain in pains:
        p = left_frame.add_paragraph()
        p.text = pain
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(12)
        p.space_after = Pt(12)

    # 右侧：我们的机会
    right_box = slide.shapes.add_textbox(
        Inches(5.3), Inches(1.5), Inches(4.2), Inches(5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    p = right_frame.paragraphs[0]
    p.text = "我们的机会"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_after = Pt(20)

    opportunities = [
        "✅ 万亿级市场\n   2025年达3000亿元\n   年复合增长率30%+",
        "✅ 政策支持\n   \"十四五\"重点扶持\n   产业升级需求",
        "✅ 技术成熟\n   AI大模型突破\n   边缘计算+5G",
        "✅ 市场空白\n   无成熟通用OS\n   先发优势明显"
    ]

    for opp in opportunities:
        p = right_frame.add_paragraph()
        p.text = opp
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(12)
        p.space_after = Pt(12)


def create_how_slide(prs):
    """创建怎么做页（技术架构）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "准备怎么做？核心技术架构"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 四层架构
    layers = [
        ("感知层 Perception", "多传感器融合\n• 激光雷达 + 视觉 + IMU\n• 点云处理 + 目标检测"),
        ("认知层 Cognitive", "World Model 世界模型\n• 环境语义理解\n• Chain-of-Thought 推理"),
        ("规划层 Planning", "HTN 分层任务规划\n• 三层规划器（任务/技能/动作）\n• 动态插入 + 失败恢复"),
        ("执行层 Execution", "自适应执行引擎\n• 实时监控 + 异常处理\n• 自动重规划")
    ]

    y_start = 1.5
    for i, (title, desc) in enumerate(layers):
        # 层次框
        layer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_start), Inches(9), Inches(1)
        )
        layer_frame = layer_box.text_frame
        layer_frame.word_wrap = True

        p = layer_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)

        p = layer_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 1

        y_start += 1.1

    # 平台支持
    platform_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.2), Inches(9), Inches(1.2)
    )
    platform_frame = platform_box.text_frame
    platform_frame.word_wrap = True
    platform_frame.text = "🚁 无人机 Drone      🚗 无人车 UGV      ⛵ 无人船 USV"
    p = platform_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


def create_advantages_slide(prs):
    """创建核心技术优势页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "核心技术优势"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 6大优势
    advantages = [
        ("✅ 统一抽象层", "一次开发，三平台复用\n代码复用率90%+", "左"),
        ("✅ World Model", "多模态传感器融合\n实时语义理解", "中"),
        ("✅ HTN智能规划", "自然语言任务分解\n自动推理决策", "右"),
        ("✅ CoT可解释AI", "Chain-of-Thought推理\n决策过程透明可追溯", "左"),
        ("✅ 自适应执行", "实时监控状态\n自动异常恢复", "中"),
        ("✅ 开放生态", "标准化接口\n支持第三方扩展", "右")
    ]

    y_positions = [1.5, 1.5, 1.5, 3.2, 3.2, 3.2]
    x_positions = [0.5, 3.3, 6.1, 0.5, 3.3, 6.1]

    for i, (title, desc, pos) in enumerate(advantages):
        box = slide.shapes.add_textbox(
            Inches(x_positions[i]), Inches(y_positions[i]), Inches(2.7), Inches(1.5)
        )
        frame = box.text_frame
        frame.word_wrap = True

        p = frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)

        p = frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # 技术壁垒数据
    data_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(5.5), Inches(9), Inches(1.5)
    )
    data_frame = data_box.text_frame
    data_frame.word_wrap = True
    data_frame.text = "🎯 技术壁垒：168项测试用例 · 100%通过 · 完整测试体系\n🔒 全栈自研：感知 + 认知 + 规划 + 执行 4大核心模块"
    p = data_frame.paragraphs[0]
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 102, 0)
    p.alignment = PP_ALIGN.CENTER


def create_drone_applications_slide(prs):
    """创建无人机应用场景页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🚁 无人机应用场景"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 应用场景
    apps = [
        ("应急救援", "搜索被困人员 · 自动标记位置\n空投物资 · 效率提升10倍"),
        ("电力巡检", "沿线自主飞行 · AI识别缺陷\n自动生成报告 · 成本1/5"),
        ("物流配送", "精准投放 · 智能避障\n10分钟送达 · 城市环境导航"),
        ("农业植保", "精准喷洒 · 病虫害检测\n节约50%农药 · 产量提升20%")
    ]

    y_start = 1.5
    for title, desc in apps:
        box = slide.shapes.add_textbox(
            Inches(1), Inches(y_start), Inches(8), Inches(1.1)
        )
        frame = box.text_frame
        frame.word_wrap = True

        p = frame.paragraphs[0]
        p.text = f"🎯 {title}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 102, 0)

        p = frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(8)

        y_start += 1.2


def create_ugv_applications_slide(prs):
    """创建无人车应用场景页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🚗 无人车应用场景"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 应用场景
    apps = [
        ("仓储物流", "多车协同 · 智能避障\n精准抓取 · 效率提升3倍"),
        ("安保巡逻", "全天候监控 · 异常检测\n人脸识别 · 成本降低70%"),
        ("环境检测", "气体泄漏监测 · 污染源定位\n自动采样 · 实时预警"),
        ("室内服务", "自主导航 · 机械臂操作\n文档整理 · 办公自动化")
    ]

    y_start = 1.5
    for title, desc in apps:
        box = slide.shapes.add_textbox(
            Inches(1), Inches(y_start), Inches(8), Inches(1.1)
        )
        frame = box.text_frame
        frame.word_wrap = True

        p = frame.paragraphs[0]
        p.text = f"🎯 {title}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 102, 0)

        p = frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(8)

        y_start += 1.2


def create_usv_applications_slide(prs):
    """创建无人船应用场景页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "⛵ 无人船应用场景"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 应用场景
    apps = [
        ("水质监测", "自主航行 · 多点采样\n实时检测 · 成本降低80%"),
        ("水下巡检", "声呐+视觉双模检测\n裂缝识别 · 3D建模"),
        ("海上救援", "多船协同 · 生命体征探测\n自动投放救生设备"),
        ("海洋科研", "鱼群跟踪 · 环境数据记录\n行为分析 · 生态保护")
    ]

    y_start = 1.5
    for title, desc in apps:
        box = slide.shapes.add_textbox(
            Inches(1), Inches(y_start), Inches(8), Inches(1.1)
        )
        frame = box.text_frame
        frame.word_wrap = True

        p = frame.paragraphs[0]
        p.text = f"🎯 {title}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 102, 0)

        p = frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(8)

        y_start += 1.2


def create_collaboration_slide(prs):
    """创建空地水协同场景页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "🌟 空地水协同场景"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 协同场景
    scenarios = [
        ("空地水一体化救援",
         "无人机搜索发现被困者 → 空投物资\n无人车快速接应转运\n无人船负责水域搜救",
         "黄金救援时间 · 挽救生命"),
        ("跨域物流配送",
         "无人机最后一公里配送\n无人车干线运输\n无人船跨水域物流",
         "全场景覆盖 · 成本优化30%"),
        ("立体环境监测",
         "无人机高空大范围扫描\n无人车地面详细检测\n无人船水域采样分析",
         "多维度数据 · 环保决策支持")
    ]

    y_start = 1.5
    for title, desc, benefit in scenarios:
        # 场景框
        scenario_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_start), Inches(9), Inches(1.4)
        )
        scenario_frame = scenario_box.text_frame
        scenario_frame.word_wrap = True

        p = scenario_frame.paragraphs[0]
        p.text = f"🎯 {title}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 102, 0)
        p.space_after = Pt(4)

        p = scenario_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(4)

        p = scenario_frame.add_paragraph()
        p.text = f"💎 {benefit}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.font.bold = True

        y_start += 1.6


def create_business_model_slide(prs):
    """创建商业模式页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "商业模式与收入预测"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 收入预测表
    table_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.3), Inches(7), Inches(2.2)
    )
    table_frame = table_box.text_frame
    table_frame.word_wrap = True

    p = table_frame.paragraphs[0]
    p.text = "📊 收入预测"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 102, 0)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    rows = [
        ("年份", "平台数量", "收入来源", "预期收入"),
        ("Y1", "1,000台", "企业授权 + 政府项目", "¥500万"),
        ("Y2", "10,000台", "企业授权 + 生态分成", "¥5000万"),
        ("Y3", "50,000台", "全线产品", "¥3亿")
    ]

    for i, (year, platforms, source, revenue) in enumerate(rows):
        p = table_frame.add_paragraph()
        p.text = f"{year}  |  {platforms}  |  {source}  |  {revenue}"
        p.font.size = Pt(16)
        if i == 0:
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 204)
        else:
            p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)

    # 商业模式
    model_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.8), Inches(9), Inches(2.5)
    )
    model_frame = model_box.text_frame
    model_frame.word_wrap = True

    p = model_frame.paragraphs[0]
    p.text = "💰 三大商业模式"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 102, 0)
    p.space_after = Pt(16)

    models = [
        ("🏢 To B - 企业授权", "机器人制造商、系统集成商", "¥5,000-50,000/台/年"),
        ("🏛️  To G - 政府项目", "应急救援、边境巡逻、环境监测", "¥100万-1000万/项目"),
        ("👨‍💻 To D - 开发者生态", "开放API + 应用商店", "API付费 + 30%抽成")
    ]

    for title, customer, pricing in models:
        p = model_frame.add_paragraph()
        p.text = f"{title}\n   客户：{customer}\n   收费：{pricing}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(12)


def create_competitive_slide(prs):
    """创建竞争优势页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "为什么选择我们？"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 竞争优势
    advantages = [
        ("vs. 传统机器人厂商",
         "❌ 他们：每款产品独立开发，周期长、成本高\n✅ 我们：统一平台，一次开发多平台复用",
         "left"),
        ("vs. ROS（Robot Operating System）",
         "❌ ROS：只是通信中间件，缺乏智能决策能力\n✅ 我们：完整的感知-认知-规划-执行闭环",
         "left"),
        ("vs. 大厂方案（如Apollo、DJI）",
         "❌ 他们：封闭生态，仅支持自家平台\n✅ 我们：开放架构，支持所有小微特平台",
         "left"),
        ("我们的核心壁垒",
         "🔒 全栈自研，拥有核心知识产权\n🎯 完整测试体系，168项测试100%通过\n🌍 开放架构，支持所有小微特平台\n♻️ 统一抽象层，一次开发三平台复用\n🧠 可解释AI，Chain-of-Thought推理",
         "left")
    ]

    y_start = 1.3
    for title, content, pos in advantages:
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_start), Inches(9), Inches(1.4)
        )
        frame = box.text_frame
        frame.word_wrap = True

        p = frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 102, 0)
        p.space_after = Pt(6)

        p = frame.add_paragraph()
        p.text = content
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)

        y_start += 1.5


def create_roadmap_slide(prs):
    """创建路线图页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "下一步计划（6个月）"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 三个阶段
    phases = [
        ("🔧 技术完善（2个月）",
         "• 完成认知层CoT推理引擎\n• 优化感知层实时性能\n• 集成3个真实平台（无人机、无人车、无人船）",
         "left"),
        ("🎬 示范应用（2个月）",
         "• 3个典型场景Demo（物流、巡检、救援）\n• 客户演示视频\n• 性能测试报告",
         "left"),
        ("💼 商业落地（2个月）",
         "• 签约3-5家意向客户\n• 完成种子轮融资\n• 组建商务团队",
         "left")
    ]

    y_start = 1.5
    for title, content, pos in phases:
        box = slide.shapes.add_textbox(
            Inches(0.5), Inches(y_start), Inches(9), Inches(1.5)
        )
        frame = box.text_frame
        frame.word_wrap = True

        p = frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.space_after = Pt(8)

        p = frame.add_paragraph()
        p.text = content
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)

        y_start += 1.6


def create_financing_slide(prs):
    """创建融资需求页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = "融资需求"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 102, 204)
    title_para.alignment = PP_ALIGN.CENTER

    # 融资额度
    amount_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.3), Inches(8), Inches(1.2)
    )
    amount_frame = amount_box.text_frame
    amount_frame.word_wrap = True
    amount_frame.text = "💰 融资需求：¥500万 - ¥1000万"
    p = amount_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 102, 0)
    p.alignment = PP_ALIGN.CENTER

    # 资金用途
    usage_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.8), Inches(9), Inches(3.5)
    )
    usage_frame = usage_box.text_frame
    usage_frame.word_wrap = True

    p = usage_frame.paragraphs[0]
    p.text = "📊 资金用途"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.space_after = Pt(20)

    usages = [
        ("🔬 技术研发 40%", "核心算法优化 · 平台集成 · 性能提升"),
        ("📢 市场推广 30%", "品牌建设 · 行业展会 · 客户获取"),
        ("👥 团队扩张 20%", "技术人才 · 商务团队 · 运营管理"),
        ("💼 备用金 10%", "风险储备 · 应急资金")
    ]

    for title, desc in usages:
        p = usage_frame.add_paragraph()
        p.text = f"{title}\n   {desc}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_after = Pt(16)


def create_contact_slide(prs):
    """创建联系页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    background = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 26, 26)
    background.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = "让每一个机器人都能理解世界，自主决策！\n\nBrain - 通用小微特机器人智能操作系统"
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 102, 204)
    p.alignment = PP_ALIGN.CENTER

    # 联系方式
    contact_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(5), Inches(7), Inches(2)
    )
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    contact_frame.text = "📧 [您的联系邮箱]\n📱 [您的联系电话]\n🏢 [您的公司地址]\n\n感谢您的关注与支持！"
    p = contact_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER


def main():
    """生成PPT"""
    prs = Presentation()

    # 设置幻灯片尺寸为16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 创建所有幻灯片
    create_title_slide(prs)
    create_why_slide(prs)
    create_how_slide(prs)
    create_advantages_slide(prs)
    create_drone_applications_slide(prs)
    create_ugv_applications_slide(prs)
    create_usv_applications_slide(prs)
    create_collaboration_slide(prs)
    create_business_model_slide(prs)
    create_competitive_slide(prs)
    create_roadmap_slide(prs)
    create_financing_slide(prs)
    create_contact_slide(prs)

    # 保存PPT
    output_path = "/media/yangyuhui/CODES1/Brain/docs/investment_pitch/Brain_Investment_Pitch.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页幻灯片")


if __name__ == "__main__":
    main()
